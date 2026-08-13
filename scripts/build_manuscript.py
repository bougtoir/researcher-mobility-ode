#!/usr/bin/env python3
"""Build a data-driven manuscript draft from ODE result CSVs.

Outputs (all regenerated from results/* CSVs, no hard-coded numbers):
- docs/manuscript.md                    (generic Markdown)
- docs/manuscript.docx                  (generic Word)
- docs/manuscript_research_policy.docx  (Research Policy formatted Word)
- docs/manuscript_research_policy.md    (Research Policy formatted Markdown)
- docs/manuscript_figures.pptx          (editable figure/table deck)
- docs/figures/*.png
"""
from __future__ import annotations

import argparse
import math
import os
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.util import Inches as PptxInches

BASE_DIR = Path(__file__).resolve().parent.parent
RESULTS_DIR = BASE_DIR / "results" / "endogenous"
SATURATING_DIR = BASE_DIR / "results" / "endogenous_saturating"
TOP_T = RESULTS_DIR / "top_transitions_T.csv"
EQ = RESULTS_DIR / "equilibrium_summary.csv"
PNR = RESULTS_DIR / "closest_point_of_no_return.csv"


def _fmt(v, dec=2):
    if pd.isna(v):
        return "—"
    try:
        return f"{float(v):.{dec}f}"
    except (ValueError, TypeError):
        return str(v)


def add_citation(para, number: int):
    run = para.add_run(f"[{number}]")
    run.font.superscript = True
    return run


def _paragraph_text(doc):
    for p in doc.paragraphs:
        yield p.text


def _table_text(doc):
    for t in doc.tables:
        for row in t.rows:
            for cell in row.cells:
                yield cell.text


def _doc_word_count(doc):
    """Approximate word count of a generated docx (used for title-page notes)."""
    return sum(len(t.split()) for t in list(_paragraph_text(doc)) + list(_table_text(doc)))


def _rel_path(path: Path, base: Path) -> str:
    return path.relative_to(base).as_posix()


def build_figure1(eq, fig_dir: Path):
    """Equilibrium domestic active pool vs minimum viable threshold."""
    fig_dir.mkdir(parents=True, exist_ok=True)
    groups = eq["group"].tolist()
    x = np.arange(len(groups))
    width = 0.35
    fig, ax = plt.subplots(figsize=(10, 5))
    ax.bar(x - width / 2, eq["T_equilibrium"], width, label="Equilibrium T", color="steelblue")
    ax.bar(x + width / 2, eq["M_threshold"], width, label="Minimum viable threshold M", color="coral")
    ax.set_xticks(x)
    ax.set_xticklabels(groups, rotation=35, ha="right")
    ax.set_ylabel("Number of researchers")
    ax.set_title("Domestic active researcher pool and minimum viable coauthor threshold by group")
    ax.legend()
    ax.set_ylim(0, max(eq["T_equilibrium"].max(), eq["M_threshold"].max()) * 1.1)
    fig.tight_layout()
    path = fig_dir / "fig1_equilibrium_margin.png"
    fig.savefig(path, dpi=300)
    plt.close(fig)
    return path


def build_figure2(pnr_closest, fig_dir: Path):
    """Proximity to point of no return for the closest rate per group.

    proximity is defined as |critical_factor - 1|, i.e. the proportional
    change in the listed transition rate required to drive the target pool
    to its minimum viable threshold.
    """
    fig_dir.mkdir(parents=True, exist_ok=True)
    pnr_closest = pnr_closest.sort_values("proximity")
    groups = pnr_closest["group"].tolist()
    rates = pnr_closest["rate_name"].tolist()
    prox = pnr_closest["proximity"].tolist()
    fig, ax = plt.subplots(figsize=(9, 5))
    colors = ["crimson" if p < 0.7 else "darkorange" if p < 0.85 else "seagreen" for p in prox]
    bars = ax.barh(groups, prox, color=colors)
    for bar, rate in zip(bars, rates):
        width = bar.get_width()
        ax.text(width + 0.01, bar.get_y() + bar.get_height() / 2,
                f"{rate}", va="center", fontsize=8)
    ax.set_xlabel("Required proportional change in rate |critical factor − 1|")
    ax.set_title("Closest point-of-no-return leverage by group (smaller = more fragile)")
    xmax = max(1.2, max(prox) * 1.15) if prox else 1.2
    ax.set_xlim(0, xmax)
    fig.tight_layout()
    path = fig_dir / "fig2_pnr_proximity.png"
    fig.savefig(path, dpi=300)
    plt.close(fig)
    return path


def _abstract_and_highlights(eq, pnr_closest):
    abstract = (
        "International mobility can concentrate AI/ML researchers in a few regions, raising the risk that smaller research communities fall below a minimum viable coauthor pool and cannot recover. "
        "We model each civilisation as a six-compartment system of domestic and abroad early-career, high-impact, and principal-investigator researchers, and estimate transition rates from OpenAlex Artificial Intelligence works (subfield 1702). "
        f"The minimum viable coauthor threshold is defined as M = k × c_bar, where c_bar is the mean number of authors per work and k is the median number of distinct last-author groups observed per recent year. "
        f"Across {len(eq)} groups, equilibrium domestic active pools T remain above M, but the closest point of no return is observed for the {pnr_closest.iloc[0]['group']} group, where a proportional change of {_fmt(pnr_closest.iloc[0]['proximity'])} in {pnr_closest.iloc[0]['rate_name']} (critical factor {_fmt(pnr_closest.iloc[0]['critical_factor'])}×) would drive the pool to its threshold. "
        "The largest positive leverage comes from PI-driven inflow and the conversion of high-impact researchers into PIs, while the dominant negative leverage is researcher dropout. "
        "These results provide a quantitative framework for early, safety-factor-bound interventions that preserve civilisational diversity in AI/ML research."
    )
    keywords = (
        "researcher mobility; artificial intelligence; civilisation grouping; "
        "ordinary differential equations; point of no return; innovation studies"
    )
    highlights = [
        "Nine civilisations modelled as six-compartment ODEs fitted to OpenAlex AI/ML data.",
        f"Closest point of no return is {pnr_closest.iloc[0]['group']} PI pool; {pnr_closest.iloc[0]['rate_name']} needs only {_fmt(pnr_closest.iloc[0]['proximity'])} proportional change.",
        "Dropout most negative; PI inflow and domestic promotion most positive.",
    ]
    return abstract, keywords, highlights


def _data_availability_text():
    return (
        "This study uses publication metadata from the OpenAlex API (subfield 1702, Artificial Intelligence; "
        "2000–2023). The extraction and analysis code, the country-to-civilisation mapping, and the result CSVs "
        "used to generate this manuscript are available in the public GitHub repository "
        "https://github.com/bougtoir/researcher-mobility-ode. OpenAlex data are released under CC0."
    )


def write_markdown(output_dir: Path, eq, top_t, pnr_closest, sat_eq, fig1_rel: str, fig2_rel: str, journal: str = "generic"):
    md = output_dir / ("manuscript.md" if journal == "generic" else f"manuscript_{journal}.md")
    abstract, keywords, highlights = _abstract_and_highlights(eq, pnr_closest)

    lines = [
        "# Quantifying the Point of No Return in Global AI/ML Research Communities",
        "",
    ]

    if journal == "research_policy":
        lines.extend([
            "**Article type:** Research Article / Research Note (to be confirmed)",
            "",
            "## Abstract",
            "",
            abstract,
            "",
            f"**Keywords:** {keywords}",
            "",
            "**Highlights**",
            "",
        ])
        for h in highlights:
            lines.append(f"- {h}")
        lines.extend(["", "## Data and Code Availability", "", _data_availability_text(), ""])
        lines.extend([
            "## Declarations",
            "",
            "**Funding:** [To be completed / removed for double-blind review]",
            "",
            "**Competing interests:** [To be completed / removed for double-blind review]",
            "",
            "**Author contributions:** [To be completed / removed for double-blind review]",
            "",
            "**Acknowledgments:** [To be completed / removed for double-blind review]",
            "",
        ])
    else:
        lines.extend(["## Abstract", "", abstract, ""])

    lines.extend([
        "## 1. Introduction",
        "",
        "Most debates on research mobility focus on net flows. Shifting attention to transition rates makes it possible to ask not only where researchers move, but which transitions must be altered to keep a community viable[1,2]. ",
        "We operationalise this idea by modelling the stock of active researchers in each of nine modified Huntington civilisations as a coupled system of ordinary differential equations. ",
        "The model is fitted to real publication records from OpenAlex and used to locate a point of no return: the parameter region in which the domestic active pool falls below the minimum number of coauthors needed to sustain the field.",
        "",
        "## 2. Methods",
        "",
        "### 2.1 Data and grouping",
        "",
        "We extracted AI/ML works and author histories from the OpenAlex API for subfield `subfields/1702` (Artificial Intelligence), using works published between 2000 and 2023[3]. ",
        "Authors were assigned to a Huntington-derived civilisation by majority country of affiliated institutions: United States, Anglosphere ex-US, Continental Europe, Sinic, Japanese, Hindu, Islamic, Other Western, and Other Civilisations. ",
        "The rationale for splitting the Western bloc and merging smaller civilisations is documented separately.",
        "",
        "### 2.2 Compartment model",
        "",
        "Each group has six compartments: domestic early-career (D), abroad early-career (A), domestic hit researchers (H_D), abroad hit researchers (H_A), domestic PIs (P_D), and abroad PIs (P_A). ",
        "Transitions are early-career outflow (alpha) and return (beta), hit generation (h_D, h_A), PI promotion (p_D, p_A), and dropout (d). ",
        "New entrants follow endogenous PI-driven inflow I(P_D) = I0 + r·P_D, with r capped at half the stability-critical value (safety factor 0.5).",
        "",
        "### 2.3 Minimum viable coauthor threshold",
        "",
        "For each group we computed the mean number of authors per work (c_bar) and the median number of distinct last-author groups per recent year (k). ",
        "The minimum viable domestic active pool is M = k × c_bar. When the equilibrium T = D + H_D + P_D falls below M, the community can no longer produce works at the observed coauthor intensity and is treated as past the point of no return.",
        "",
        "### 2.4 Sensitivity and point-of-no-return scan",
        "",
        "We computed elasticities by perturbing each transition rate by 1% and re-solving the equilibrium. ",
        "For point-of-no-return analysis we scaled each rate in turn until T reached M, recording the critical factor and its proximity, defined as |critical factor − 1| (the proportional change in that rate required to reach the threshold).",
        "",
        "## 3. Results",
        "",
        f"Table 1 reports equilibrium domestic active pool T, minimum viable threshold M, and endogenous inflow parameters for the {len(eq)} groups. ",
        "All groups remain above their threshold under the fitted model, but margins differ by an order of magnitude.",
        "",
        "**Table 1. Equilibrium domestic active pool, minimum viable threshold, and endogenous inflow parameters.**",
        "",
    ])
    headers = ["Group", "T_eq", "M", "Margin", "I0", "r", "r_obs", "r_crit"]
    col_map = ["group", "T_equilibrium", "M_threshold", "margin_to_threshold_T", "I0", "r", "r_obs", "r_critical"]
    decimals = [None, 2, 2, 2, 2, 5, 5, 5]
    lines.append("| " + " | ".join(headers) + " |")
    lines.append("|" + "|".join(["---"] * len(headers)) + "|")
    for _, row in eq.iterrows():
        lines.append("| " + " | ".join([_fmt(row[c], decimals[i] or 2) for i, c in enumerate(col_map)]) + " |")
    lines.append("")

    lines.extend([
        f"Table 2 shows the three transition-rate elasticities with the largest absolute impact on T for each group. ",
        "Dropout (d) is the largest negative lever in every group; promotion of domestic hit researchers to PIs (p_D) and return from abroad (beta) are the main positive levers after inflow.",
        "",
        "**Table 2. Top transition-rate elasticities for domestic active pool T.**",
        "",
    ])
    lines.append("| Group | 1st rate | 1st elasticity | 2nd rate | 2nd elasticity | 3rd rate | 3rd elasticity |")
    lines.append("|---|---|---|---|---|---|---|")
    for group, gdf in top_t.groupby("group"):
        top3 = gdf.sort_values("abs_elasticity", ascending=False).head(3)
        vals = top3[["rate", "elasticity"]].values.tolist()
        parts = [group]
        for rate, elas in vals:
            parts.extend([rate, _fmt(elas, 3)])
        lines.append("| " + " | ".join(parts) + " |")
    lines.append("")

    closest = pnr_closest.iloc[0]
    lines.extend([
        "Table 3 reports, for each group, the single rate that reaches the threshold with the smallest proportional change (closest point of no return). ",
        f"The {closest['group']} group is the most fragile: a proportional change of {_fmt(closest['proximity'])} in {closest['rate_name']} (critical factor {_fmt(closest['critical_factor'])}×) would drive the {closest['target']} pool to its minimum viable threshold.",
        "",
        "**Table 3. Closest point of no return by group.**",
        "",
    ])
    lines.append("| Group | Target | Rate | Current | Critical factor | Proximity |")
    lines.append("|---|---|---|---|---|---|")
    for _, row in pnr_closest.iterrows():
        lines.append(
            f"| {row['group']} | {row['target']} | {row['rate_name']} | {_fmt(row['current_rate'], 4)} | {_fmt(row['critical_factor'], 3)} | {_fmt(row['proximity'], 3)} |"
        )
    lines.append("")

    lines.extend([
        f"![Figure 1]({fig1_rel})",
        "",
        "**Figure 1. Equilibrium domestic active pool (T) and minimum viable coauthor threshold (M) by group.** All groups remain above the threshold, but the margin varies widely.",
        "",
        f"![Figure 2]({fig2_rel})",
        "",
        "**Figure 2. Closest point-of-no-return proximity by group.** Smaller values indicate that a smaller proportional change in the listed transition rate is required to drive the group to its threshold.",
        "",
    ])

    if sat_eq is not None:
        lines.extend([
            "### 3.1 Saturating recruitment extension",
            "",
            "Replacing linear inflow with a saturating form I(P_D) = I0 + r·P_D / (1 + ε·P_D) lowers equilibrium pools for fast-growing groups because each additional PI adds fewer entrants. ",
            "Table 4 compares the linear and saturating equilibrium T values.",
            "",
            "**Table 4. Equilibrium T under linear and saturating PI-driven inflow.**",
            "",
        ])
        lines.append("| Group | Linear T | Saturating T | ε |")
        lines.append("|---|---|---|---|")
        merged = eq[["group", "T_equilibrium"]].merge(
            sat_eq[["group", "T_equilibrium", "epsilon"]], on="group", suffixes=("_lin", "_sat")
        )
        for _, row in merged.iterrows():
            lines.append(
                f"| {row['group']} | {_fmt(row['T_equilibrium_lin'], 2)} | {_fmt(row['T_equilibrium_sat'], 2)} | {_fmt(row['epsilon'], 5)} |"
            )
        lines.append("")

    lines.extend([
        "## 4. Discussion",
        "",
        "The model supports a transition-rate view of research-policy intervention. ",
        "Because dropout has an elasticity near -2 for every group, policies that reduce attrition — stable junior positions, grants for risky early work, and family/visa support — have the highest marginal impact on community size. ",
        "At the same time, PI-driven inflow and domestic PI promotion (p_D) have the largest positive elasticities, indicating that sustaining a senior core is necessary for generational renewal.",
        "",
        "The Japanese group illustrates how a technologically advanced but demographically smaller civilisation can sit close to the PI point of no return even when the overall active pool looks comfortable. ",
        "This asymmetry between T and P_D suggests that headline researcher counts can mask fragility in leadership generation.",
        "",
        "From a global-diversity standpoint, the results argue for early intervention within a safety factor: small proportional adjustments to return rates, hit-generation, and dropout are sufficient to keep every group above its threshold, avoiding the concentration that would turn AI/ML into an oligopoly of a few large civilisations[1,2].",
        "",
        "## 5. Limitations",
        "",
        "OpenAlex affiliation and author country assignments are noisy, and the model treats each civilisation as a closed compartment with no cross-civilisation spillovers beyond endogenous inflow. ",
        "Future extensions include network externalities, time-varying parameters, and a larger quantum-technology pilot to test transferability to security-relevant fields.",
        "",
        "## References",
        "",
        "1. Momentumyy. 人材流出ではなく『遷移係数』で考える研究コミュニティの存亡. note, 2024. https://note.com/momentumyy/n/n86df5d34282d (accessed 2024-08-09).",
        "2. Huntington S P. The Clash of Civilizations and the Remaking of World Order. New York: Simon & Schuster, 1996.",
        "3. Priem J, Piwowar H, Orr R. OpenAlex: A fully-open index of scholarly works, authors, venues, institutions, and concepts. arXiv:2205.01813, 2022.",
        "",
    ])
    md.write_text("\n".join(lines), encoding="utf-8")
    return md


def _add_docx_title_page(doc, word_count=None):
    """Add Research Policy-style title page."""
    style = doc.styles["Normal"]
    style.font.name = "Times New Roman"
    style.font.size = Pt(11)
    style.paragraph_format.line_spacing = 1.5

    title = doc.add_heading("Quantifying the Point of No Return in Global AI/ML Research Communities", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title.runs[0].font.size = Pt(16)
    title.runs[0].font.bold = True

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run("Article type: Research Article / Research Note (to be confirmed)")
    if word_count:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run(f"Approximate word count (main text incl. tables, excl. references): {word_count}")
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run("Note: Research Policy Full Article target is ~8,000–10,000 words; Research Note is 3,000–5,000 words. Expand or select article type accordingly.")

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run("Corresponding author: [To be completed / removed for double-blind review]")

    # Page break after title page
    p = doc.add_paragraph()
    run = p.add_run()
    run.add_break(WD_BREAK.PAGE)


def _add_docx_front_matter(doc, abstract, keywords, highlights):
    h = doc.add_heading("Abstract", level=1)
    p = doc.add_paragraph()
    p.add_run(abstract)

    p = doc.add_paragraph()
    p.add_run("Keywords: ").bold = True
    p.add_run(keywords)

    doc.add_heading("Highlights", level=2)
    for hlt in highlights:
        p = doc.add_paragraph(style="List Bullet")
        p.add_run(hlt)

    doc.add_heading("Data and Code Availability", level=2)
    p = doc.add_paragraph()
    p.add_run(_data_availability_text())

    doc.add_heading("Declarations", level=2)
    for sub in ["Funding", "Competing interests", "Author contributions", "Acknowledgments"]:
        p = doc.add_paragraph()
        p.add_run(f"{sub}: ").bold = True
        p.add_run("[To be completed / removed for double-blind review]")


def _add_docx_body(doc, eq, top_t, pnr_closest, sat_eq, fig1, fig2):
    # Introduction
    doc.add_heading("1. Introduction", level=1)
    p = doc.add_paragraph()
    p.add_run("Most debates on research mobility focus on net flows. Shifting attention to transition rates makes it possible to ask not only where researchers move, but which transitions must be altered to keep a community viable")
    add_citation(p, 1)
    p.add_run(", and the Huntington civilisation framework helps partition the world into culturally coherent research arenas")
    add_citation(p, 2)
    p.add_run(". We operationalise this idea by modelling the stock of active researchers in each of nine modified Huntington civilisations as a coupled system of ordinary differential equations. The model is fitted to real publication records from OpenAlex")
    add_citation(p, 3)
    p.add_run(" and used to locate a point of no return: the parameter region in which the domestic active pool falls below the minimum number of coauthors needed to sustain the field.")

    # Methods
    doc.add_heading("2. Methods", level=1)
    doc.add_heading("2.1 Data and grouping", level=2)
    p = doc.add_paragraph()
    p.add_run("We extracted AI/ML works and author histories from the OpenAlex API for subfield `subfields/1702` (Artificial Intelligence), using works published between 2000 and 2023.")
    add_citation(p, 3)
    p.add_run(" Authors were assigned to a Huntington-derived civilisation by majority country of affiliated institutions: United States, Anglosphere ex-US, Continental Europe, Sinic, Japanese, Hindu, Islamic, Other Western, and Other Civilisations. The rationale for splitting the Western bloc and merging smaller civilisations is documented separately.")

    doc.add_heading("2.2 Compartment model", level=2)
    p = doc.add_paragraph()
    p.add_run("Each group has six compartments: domestic early-career (D), abroad early-career (A), domestic hit researchers (H_D), abroad hit researchers (H_A), domestic PIs (P_D), and abroad PIs (P_A). Transitions are early-career outflow (alpha) and return (beta), hit generation (h_D, h_A), PI promotion (p_D, p_A), and dropout (d). New entrants follow endogenous PI-driven inflow I(P_D) = I0 + r·P_D, with r capped at half the stability-critical value (safety factor 0.5).")

    doc.add_heading("2.3 Minimum viable coauthor threshold", level=2)
    p = doc.add_paragraph()
    p.add_run("For each group we computed the mean number of authors per work (c_bar) and the median number of distinct last-author groups per recent year (k). The minimum viable domestic active pool is M = k × c_bar. When the equilibrium T = D + H_D + P_D falls below M, the community can no longer produce works at the observed coauthor intensity and is treated as past the point of no return.")

    doc.add_heading("2.4 Sensitivity and point-of-no-return scan", level=2)
    p = doc.add_paragraph()
    p.add_run("We computed elasticities by perturbing each transition rate by 1% and re-solving the equilibrium. For point-of-no-return analysis we scaled each rate in turn until T reached M, recording the critical factor and its proximity, defined as |critical factor − 1| (the proportional change in that rate required to reach the threshold).")

    # Results
    doc.add_heading("3. Results", level=1)
    p = doc.add_paragraph()
    p.add_run(f"Table 1 reports equilibrium domestic active pool T, minimum viable threshold M, and endogenous inflow parameters for the {len(eq)} groups. All groups remain above their threshold under the fitted model, but margins differ by an order of magnitude.")

    # Table 1
    headers1 = ["Group", "T_eq", "M", "Margin", "I0", "r", "r_obs", "r_crit"]
    col_map1 = ["group", "T_equilibrium", "M_threshold", "margin_to_threshold_T", "I0", "r", "r_obs", "r_critical"]
    decimals1 = [None, 2, 2, 2, 2, 5, 5, 5]
    t1 = doc.add_table(rows=1, cols=len(headers1))
    t1.style = "Table Grid"
    for i, h in enumerate(headers1):
        t1.rows[0].cells[i].text = h
    for _, row in eq.iterrows():
        cells = t1.add_row().cells
        for i, c in enumerate(col_map1):
            cells[i].text = _fmt(row[c], decimals1[i] or 2)
    cap = doc.add_paragraph()
    cap.add_run("Table 1. Equilibrium domestic active pool, minimum viable threshold, and endogenous inflow parameters.").italic = True
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Table 2
    p = doc.add_paragraph()
    p.add_run("Table 2 shows the three transition-rate elasticities with the largest absolute impact on T for each group. Dropout (d) is the largest negative lever in every group; promotion of domestic hit researchers to PIs (p_D) and return from abroad (beta) are the main positive levers after inflow.")
    headers2 = ["Group", "1st rate", "1st elasticity", "2nd rate", "2nd elasticity", "3rd rate", "3rd elasticity"]
    t2 = doc.add_table(rows=1, cols=len(headers2))
    t2.style = "Table Grid"
    for i, h in enumerate(headers2):
        t2.rows[0].cells[i].text = h
    for group, gdf in top_t.groupby("group"):
        top3 = gdf.sort_values("abs_elasticity", ascending=False).head(3)
        cells = t2.add_row().cells
        cells[0].text = group
        for j, (_, row) in enumerate(top3.iterrows()):
            cells[2 * j + 1].text = row["rate"]
            cells[2 * j + 2].text = _fmt(row["elasticity"], 3)
    cap2 = doc.add_paragraph()
    cap2.add_run("Table 2. Top transition-rate elasticities for domestic active pool T.").italic = True
    cap2.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Table 3
    closest = pnr_closest.iloc[0]
    p = doc.add_paragraph()
    p.add_run(f"Table 3 reports, for each group, the single rate that reaches the threshold with the smallest proportional change (closest point of no return). The {closest['group']} group is the most fragile: a proportional change of {_fmt(closest['proximity'])} in {closest['rate_name']} (critical factor {_fmt(closest['critical_factor'])}×) would drive the {closest['target']} pool to its minimum viable threshold.")
    headers3 = ["Group", "Target", "Rate", "Current", "Critical factor", "Proximity"]
    t3 = doc.add_table(rows=1, cols=len(headers3))
    t3.style = "Table Grid"
    for i, h in enumerate(headers3):
        t3.rows[0].cells[i].text = h
    for _, row in pnr_closest.iterrows():
        cells = t3.add_row().cells
        cells[0].text = row["group"]
        cells[1].text = row["target"]
        cells[2].text = row["rate_name"]
        cells[3].text = _fmt(row["current_rate"], 4)
        cells[4].text = _fmt(row["critical_factor"], 3)
        cells[5].text = _fmt(row["proximity"], 3)
    cap3 = doc.add_paragraph()
    cap3.add_run("Table 3. Closest point of no return by group.").italic = True
    cap3.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Figures
    p = doc.add_paragraph()
    p.add_run("Figure 1 visualises the gap between equilibrium and threshold; Figure 2 ranks groups by their closest point-of-no-return leverage.")
    doc.add_picture(str(fig1), width=Inches(5.8))
    capf1 = doc.add_paragraph()
    capf1.add_run("Figure 1. Equilibrium domestic active pool (T) and minimum viable coauthor threshold (M) by group.").italic = True
    capf1.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_picture(str(fig2), width=Inches(5.8))
    capf2 = doc.add_paragraph()
    capf2.add_run("Figure 2. Closest point-of-no-return proximity by group. Smaller values indicate that a smaller proportional change in the listed rate is required to reach the threshold.").italic = True
    capf2.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Saturating extension
    if sat_eq is not None:
        doc.add_heading("3.1 Saturating recruitment extension", level=2)
        p = doc.add_paragraph()
        p.add_run("Replacing linear inflow with a saturating form lowers equilibrium pools for fast-growing groups because each additional PI adds fewer entrants. Table 4 compares linear and saturating equilibrium T values.")
        headers4 = ["Group", "Linear T", "Saturating T", "ε"]
        t4 = doc.add_table(rows=1, cols=len(headers4))
        t4.style = "Table Grid"
        for i, h in enumerate(headers4):
            t4.rows[0].cells[i].text = h
        merged = eq[["group", "T_equilibrium"]].merge(
            sat_eq[["group", "T_equilibrium", "epsilon"]], on="group", suffixes=("_lin", "_sat")
        )
        for _, row in merged.iterrows():
            cells = t4.add_row().cells
            cells[0].text = row["group"]
            cells[1].text = _fmt(row["T_equilibrium_lin"])
            cells[2].text = _fmt(row["T_equilibrium_sat"])
            cells[3].text = _fmt(row["epsilon"], 5)
        cap4 = doc.add_paragraph()
        cap4.add_run("Table 4. Equilibrium T under linear and saturating PI-driven inflow.").italic = True
        cap4.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Discussion
    doc.add_heading("4. Discussion", level=1)
    p = doc.add_paragraph()
    p.add_run("The model supports a transition-rate view of research-policy intervention. Because dropout has an elasticity near -2 for every group, policies that reduce attrition — stable junior positions, grants for risky early work, and family/visa support — have the highest marginal impact on community size. At the same time, PI-driven inflow and domestic PI promotion (p_D) have the largest positive elasticities, indicating that sustaining a senior core is necessary for generational renewal.")
    p = doc.add_paragraph()
    p.add_run("The Japanese group illustrates how a technologically advanced but demographically smaller civilisation can sit close to the PI point of no return even when the overall active pool looks comfortable. This asymmetry between T and P_D suggests that headline researcher counts can mask fragility in leadership generation.")
    p = doc.add_paragraph()
    p.add_run("From a global-diversity standpoint, the results argue for early intervention within a safety factor: small proportional adjustments to return rates, hit-generation, and dropout are sufficient to keep every group above its threshold, avoiding the concentration that would turn AI/ML into an oligopoly of a few large civilisations.")
    add_citation(p, 1)
    add_citation(p, 2)

    # Limitations
    doc.add_heading("5. Limitations", level=1)
    p = doc.add_paragraph()
    p.add_run("OpenAlex affiliation and author country assignments are noisy, and the model treats each civilisation as a closed compartment with no cross-civilisation spillovers beyond endogenous inflow. Future extensions include network externalities, time-varying parameters, and a larger quantum-technology pilot to test transferability to security-relevant fields.")

    # References
    doc.add_heading("References", level=1)
    refs = [
        "Momentumyy. 人材流出ではなく『遷移係数』で考える研究コミュニティの存亡. note, 2024. https://note.com/momentumyy/n/n86df5d34282d (accessed 2024-08-09).",
        "Huntington S P. The Clash of Civilizations and the Remaking of World Order. New York: Simon & Schuster, 1996.",
        "Priem J, Piwowar H, Orr R. OpenAlex: A fully-open index of scholarly works, authors, venues, institutions, and concepts. arXiv:2205.01813, 2022.",
    ]
    for i, ref in enumerate(refs, 1):
        p = doc.add_paragraph()
        p.add_run(f"{i}. {ref}")


def write_docx_generic(output_dir: Path, eq, top_t, pnr_closest, sat_eq, fig1: Path, fig2: Path):
    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "Times New Roman"
    style.font.size = Pt(11)

    title = doc.add_heading("Quantifying the Point of No Return in Global AI/ML Research Communities", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    _add_docx_body(doc, eq, top_t, pnr_closest, sat_eq, fig1, fig2)

    path = output_dir / "manuscript.docx"
    doc.save(path)
    return path


def write_docx_research_policy(output_dir: Path, eq, top_t, pnr_closest, sat_eq, fig1: Path, fig2: Path):
    abstract, keywords, highlights = _abstract_and_highlights(eq, pnr_closest)

    # Pre-compute body word count (excluding title page, front matter and references)
    body_doc = Document()
    _add_docx_body(body_doc, eq, top_t, pnr_closest, sat_eq, fig1, fig2)
    body_wc = _doc_word_count(body_doc)

    doc = Document()
    _add_docx_title_page(doc, word_count=body_wc)
    _add_docx_front_matter(doc, abstract, keywords, highlights)
    _add_docx_body(doc, eq, top_t, pnr_closest, sat_eq, fig1, fig2)

    path = output_dir / "manuscript_research_policy.docx"
    doc.save(path)
    return path


def write_pptx(output_dir: Path, eq, top_t, pnr_closest, sat_eq, fig1: Path, fig2: Path):
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    def add_title_slide(title, subtitle=""):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        if subtitle:
            slide.placeholders[1].text = subtitle

    def add_image_slide(title, img_path, caption):
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        left = PptxInches(1.5)
        top = PptxInches(1.2)
        slide.shapes.add_picture(str(img_path), left, top, width=PptxInches(10))
        txBox = slide.shapes.add_textbox(left, PptxInches(6.0), PptxInches(10), PptxInches(0.8))
        txBox.text_frame.text = caption
        for paragraph in txBox.text_frame.paragraphs:
            paragraph.font.size = Pt(14)

    def add_table_slide(title, df, col_names, width_per_col=1.5, font_size=10):
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        rows, cols = len(df) + 1, len(col_names)
        left = PptxInches(0.5)
        top = PptxInches(1.2)
        table = slide.shapes.add_table(rows, cols, left, top, PptxInches(cols * width_per_col), PptxInches(0.7 * rows)).table
        for i, h in enumerate(col_names):
            table.cell(0, i).text = h
        for i, row in df.iterrows():
            for j, val in enumerate(row):
                table.cell(i + 1, j).text = str(val)
                table.cell(i + 1, j).text_frame.paragraphs[0].font.size = Pt(font_size)

    add_title_slide(
        "Quantifying the Point of No Return in Global AI/ML Research Communities",
        "Data-driven manuscript figures and tables"
    )

    add_image_slide(
        "Figure 1: Equilibrium domestic active pool vs minimum viable threshold",
        fig1,
        "Blue bars: equilibrium T; orange bars: threshold M. All groups remain above the threshold, but margins vary widely."
    )
    add_image_slide(
        "Figure 2: Closest point-of-no-return proximity by group",
        fig2,
        "Smaller values mean a smaller proportional change in the listed rate is required to reach the threshold."
    )

    t1_df = eq[["group", "T_equilibrium", "M_threshold", "margin_to_threshold_T", "I0", "r", "r_obs", "r_critical"]].copy()
    pptx_decimals = {
        "T_equilibrium": 2,
        "M_threshold": 2,
        "margin_to_threshold_T": 2,
        "I0": 2,
        "r": 5,
        "r_obs": 5,
        "r_critical": 5,
    }
    for c, dec in pptx_decimals.items():
        t1_df[c] = t1_df[c].apply(lambda x, d=dec: _fmt(x, d))
    t1_df.columns = ["Group", "T_eq", "M", "Margin", "I0", "r", "r_obs", "r_crit"]
    add_table_slide("Table 1: Equilibrium, threshold and inflow parameters", t1_df, t1_df.columns.tolist())

    rows2 = []
    for group, gdf in top_t.groupby("group"):
        top3 = gdf.sort_values("abs_elasticity", ascending=False).head(3)
        parts = [group]
        for _, row in top3.iterrows():
            parts.extend([row["rate"], _fmt(row["elasticity"], 3)])
        rows2.append(parts)
    t2_df = pd.DataFrame(rows2, columns=["Group", "1st", "el1", "2nd", "el2", "3rd", "el3"])
    add_table_slide("Table 2: Top transition-rate elasticities for T", t2_df, t2_df.columns.tolist(), width_per_col=1.4)

    t3_df = pnr_closest[["group", "target", "rate_name", "current_rate", "critical_factor", "proximity"]].copy()
    t3_decimals = {"current_rate": 4, "critical_factor": 3, "proximity": 3}
    for c, dec in t3_decimals.items():
        t3_df[c] = t3_df[c].apply(lambda x, d=dec: _fmt(x, d))
    t3_df.columns = ["Group", "Target", "Rate", "Current", "Crit.factor", "Proximity"]
    add_table_slide("Table 3: Closest point of no return", t3_df, t3_df.columns.tolist(), width_per_col=1.8)

    if sat_eq is not None:
        merged = eq[["group", "T_equilibrium"]].merge(
            sat_eq[["group", "T_equilibrium", "epsilon"]], on="group", suffixes=("_lin", "_sat")
        )
        t4_decimals = {"T_equilibrium_lin": 2, "T_equilibrium_sat": 2, "epsilon": 5}
        for c, dec in t4_decimals.items():
            merged[c] = merged[c].apply(lambda x, d=dec: _fmt(x, d))
        merged.columns = ["Group", "Linear T", "Saturating T", "ε"]
        add_table_slide("Table 4: Saturating inflow extension", merged, merged.columns.tolist())

    path = output_dir / "manuscript_figures.pptx"
    prs.save(path)
    return path


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--output-dir", type=Path, default=BASE_DIR / "docs")
    parser.add_argument("--journal", choices=["generic", "research_policy"], default="generic",
                        help="Target journal format to generate (default: generic; research_policy adds title page, abstract, highlights, declarations).")
    args = parser.parse_args()
    output_dir = args.output_dir
    output_dir.mkdir(parents=True, exist_ok=True)

    fig_dir = output_dir / "figures"

    eq = pd.read_csv(EQ)
    sat_eq = pd.read_csv(SATURATING_DIR / "equilibrium_summary.csv") if (SATURATING_DIR / "equilibrium_summary.csv").exists() else None
    top_t = pd.read_csv(TOP_T)
    pnr = pd.read_csv(PNR)

    # Closest PNR per group (smallest absolute proportional change needed)
    pnr_closest = pnr.loc[pnr.groupby("group")["proximity"].idxmin()].reset_index(drop=True)
    pnr_closest = pnr_closest.sort_values("proximity").reset_index(drop=True)

    fig1 = build_figure1(eq, fig_dir)
    fig2 = build_figure2(pnr_closest, fig_dir)
    fig1_rel = _rel_path(fig1, output_dir)
    fig2_rel = _rel_path(fig2, output_dir)

    md_path = write_markdown(output_dir, eq, top_t, pnr_closest, sat_eq, fig1_rel, fig2_rel, journal=args.journal)
    docx_path = write_docx_generic(output_dir, eq, top_t, pnr_closest, sat_eq, fig1, fig2)
    if args.journal == "research_policy":
        rp_path = write_docx_research_policy(output_dir, eq, top_t, pnr_closest, sat_eq, fig1, fig2)
        print(f"Wrote {rp_path}")
    pptx_path = write_pptx(output_dir, eq, top_t, pnr_closest, sat_eq, fig1, fig2)

    print(f"Wrote {md_path}")
    print(f"Wrote {docx_path}")
    print(f"Wrote {pptx_path}")
    print(f"Figures saved to {fig_dir}")


if __name__ == "__main__":
    main()

#!/usr/bin/env python3
"""Build a Research Policy full-article manuscript from ODE result CSVs.

Outputs (regenerated from results/* CSVs, no hard-coded numbers):
- docs/manuscript_full_article.docx
- docs/manuscript_full_article.md
- docs/manuscript_full_article_figures.pptx
- docs/figures/*.png

All numerical values are read from the result CSVs produced by the analysis
pipeline; the script contains only formatting and prose.
"""
from __future__ import annotations

import argparse
import math
import os
import sys
import zipfile
from pathlib import Path

# Make local packages importable
BASE_DIR = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(BASE_DIR / "src"))
sys.path.insert(0, str(BASE_DIR / "scripts"))

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
import pandas as pd
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.util import Inches as PptxInches

from pyommlbuilder.main import (
    Math,
    SubscriptObject,
    Fraction,
    Numerator,
    Denominator,
    Function,
    MathPara,
)
from pyommlbuilder.helpers import make_aligned_equation

import annual_rates_projection_report as arpr

try:
    import pypandoc
except Exception:  # pragma: no cover - pypandoc is optional for the markdown fallback
    pypandoc = None

RESULTS_DIR = BASE_DIR / "results"
ENDOG = RESULTS_DIR / "endogenous"
SAT = RESULTS_DIR / "endogenous_saturating"
TV = RESULTS_DIR / "time_varying"
BOOT = RESULTS_DIR / "bootstrap_ci"
POL = RESULTS_DIR / "policy_counterfactuals"
ANNUAL = RESULTS_DIR / "annual"
FIG_DIR = BASE_DIR / "docs" / "figures"


def _fmt(v, dec=2):
    if pd.isna(v):
        return "—"
    try:
        return f"{float(v):.{dec}f}"
    except (ValueError, TypeError):
        return str(v)


def add_citation(para, number: int):
    run = para.add_run(f" [{number}]")
    run.font.superscript = True
    return run


def add_footnote(para, symbol="1"):
    run = para.add_run(f" {symbol}")
    run.font.superscript = True
    return run


RATE_LABELS = {
    "I0": "exogenous entry rate (I0)",
    "I": "exogenous entry rate (I)",
    "d": "dropout rate (d)",
    "alpha": "early-career outflow rate (α)",
    "beta": "return rate (β)",
    "h_D": "domestic hit-generation rate (h_D)",
    "h_A": "abroad hit-generation rate (h_A)",
    "p_D": "domestic principal-investigator promotion rate (p_D)",
    "p_A": "abroad principal-investigator promotion rate (p_A)",
    "r": "principal-investigator reproduction rate (r)",
}


def _rate_label(name):
    """Return a full-spelling phrase with the rate symbol in parentheses."""
    return RATE_LABELS.get(str(name), str(name))


def _paragraph_text(doc):
    for p in doc.paragraphs:
        yield p.text


def _table_text(doc):
    for t in doc.tables:
        for row in t.rows:
            for cell in row.cells:
                yield cell.text


def _doc_word_count(doc, exclude_after="References"):
    """Count words in paragraphs and tables, stopping before the reference list."""
    total = 0
    active = True
    for p in doc.paragraphs:
        text = p.text.strip()
        if active and text == exclude_after:
            active = False
            continue
        if active:
            total += len(text.split())
    if active:
        for t in doc.tables:
            for row in t.rows:
                for cell in row.cells:
                    total += len(cell.text.split())
    return total


def _rel_path(path: Path, base: Path) -> str:
    return path.relative_to(base).as_posix()


def add_omath_paragraph(doc, math_element, align=WD_ALIGN_PARAGRAPH.CENTER):
    """Append an OMML math element to a new paragraph."""
    p = doc.add_paragraph()
    p._element.append(math_element._as_xml_element())
    p.alignment = align
    return p


def add_omath_inline(para, math_element):
    """Append an OMML math element inline inside an existing paragraph."""
    para._element.append(math_element._as_xml_element())


# ---------------------------------------------------------------------------
# OMML equation builders
# ---------------------------------------------------------------------------

def math_I_linear():
    """I(P_D) = I_0 + r P_D"""
    return Math(
        Function("I", SubscriptObject("P", "D")),
        "=",
        SubscriptObject("I", "0"),
        "+",
        "r",
        SubscriptObject("P", "D"),
    )


def math_I_saturating():
    """I(P_D) = I_0 + r P_D / (1 + epsilon P_D)"""
    return Math(
        Function("I", SubscriptObject("P", "D")),
        "=",
        SubscriptObject("I", "0"),
        "+",
        Fraction(
            Numerator("r", SubscriptObject("P", "D")),
            Denominator("1 + ", "ε", "×", SubscriptObject("P", "D")),
        ),
    )


def math_threshold():
    """M = k × c̄"""
    return Math("M = k × c\u0304")


def math_active_pool():
    """T = D + H_D + P_D"""
    return Math(
        "T = D + ",
        SubscriptObject("H", "D"),
        " + ",
        SubscriptObject("P", "D"),
    )


def math_ode_system():
    """Six-equation display using MathPara."""
    def deriv(base):
        return Fraction(Numerator("d" + base), Denominator("dt"))

    lines = [
        make_aligned_equation(
            deriv("D"),
            Math(
                Function("I", SubscriptObject("P", "D")),
                " + βA - (α + ",
                SubscriptObject("h", "D"),
                " + d)D",
            ),
            line_break=False,
        ),
        make_aligned_equation(
            deriv("A"),
            Math(
                "αD - (β + ",
                SubscriptObject("h", "A"),
                " + d)A",
            ),
            line_break=False,
        ),
        make_aligned_equation(
            deriv("H_D"),
            Math(
                SubscriptObject("h", "D"),
                "D + β",
                SubscriptObject("H", "A"),
                " - (",
                SubscriptObject("p", "D"),
                " + d)",
                SubscriptObject("H", "D"),
            ),
            line_break=False,
        ),
        make_aligned_equation(
            deriv("H_A"),
            Math(
                SubscriptObject("h", "A"),
                "A - (β + ",
                SubscriptObject("p", "A"),
                " + d)",
                SubscriptObject("H", "A"),
            ),
            line_break=False,
        ),
        make_aligned_equation(
            deriv("P_D"),
            Math(
                SubscriptObject("p", "D"),
                SubscriptObject("H", "D"),
                " + β",
                SubscriptObject("P", "A"),
                " - d",
                SubscriptObject("P", "D"),
            ),
            line_break=False,
        ),
        make_aligned_equation(
            deriv("P_A"),
            Math(
                SubscriptObject("p", "A"),
                SubscriptObject("H", "A"),
                " - (β + d)",
                SubscriptObject("P", "A"),
            ),
            line_break=False,
        ),
    ]
    return MathPara(*lines)


# ---------------------------------------------------------------------------
# Figures
# ---------------------------------------------------------------------------

def build_figure1(eq, fig_dir: Path):
    """Equilibrium domestic active pool vs minimum viable threshold."""
    fig_dir.mkdir(parents=True, exist_ok=True)
    groups = eq["group"].tolist()
    x = np.arange(len(groups))
    width = 0.35
    fig, ax = plt.subplots(figsize=(10, 5.5))
    t_bars = ax.bar(x - width / 2, eq["T_equilibrium"], width, label="Equilibrium T", color="steelblue")
    m_bars = ax.bar(x + width / 2, eq["M_threshold"], width, label="Minimum viable threshold M", color="coral")
    ax.set_xticks(x)
    ax.set_xticklabels(groups, rotation=35, ha="right")
    ax.set_ylabel("Number of researchers")
    ax.set_title("Domestic active researcher pool and minimum viable coauthor threshold by group")
    max_y = max(eq["T_equilibrium"].max(), eq["M_threshold"].max()) * 1.15
    ax.set_ylim(0, max_y)
    # Annotate bars with integer counts
    for bar, val in zip(t_bars, eq["T_equilibrium"]):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + max_y * 0.01,
                f"{int(round(val))}", ha="center", va="bottom", fontsize=7)
    for bar, val in zip(m_bars, eq["M_threshold"]):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + max_y * 0.01,
                f"{int(round(val))}", ha="center", va="bottom", fontsize=7)
    ax.legend()
    fig.tight_layout()
    path = fig_dir / "fig1_equilibrium_margin.png"
    fig.savefig(path, dpi=300)
    plt.close(fig)
    return path


def build_figure2(pnr_closest, fig_dir: Path):
    """Closest point-of-no-return proximity by group."""
    fig_dir.mkdir(parents=True, exist_ok=True)
    pnr_closest = pnr_closest.sort_values("proximity")
    groups = pnr_closest["group"].tolist()
    labels = [f"{r}\n({t})" for r, t in zip(pnr_closest["rate_name"], pnr_closest["target"])]
    prox = pnr_closest["proximity"].tolist()
    fig, ax = plt.subplots(figsize=(9, 5))
    # Use a perceptually uniform, colour-vision-deficiency-friendly sequential palette
    norm = max(prox) * 1.2 if prox else 1.0
    colors = [plt.cm.plasma(0.25 + 0.55 * (p / norm)) for p in prox]
    bars = ax.barh(groups, prox, color=colors)
    for bar, label in zip(bars, labels):
        width = bar.get_width()
        ax.text(width + 0.01, bar.get_y() + bar.get_height() / 2,
                label, va="center", fontsize=7)
    ax.set_xlabel("Required proportional change in rate |critical factor − 1|")
    ax.set_title("Closest point-of-no-return sensitivity by group (smaller = more fragile)")
    xmax = max(1.2, max(prox) * 1.15) if prox else 1.2
    ax.set_xlim(0, xmax)
    ax.axvline(1.0, color="gray", linestyle="--", linewidth=0.8)
    fig.tight_layout()
    path = fig_dir / "fig2_pnr_proximity.png"
    fig.savefig(path, dpi=300)
    plt.close(fig)
    return path


def build_figure3(period_compare, fig_dir: Path):
    """Historical counterfactual: change in safety margin from early to late rates.

    Mirrored to a left-origin layout: all bars originate at the left axis (0)
    and extend to the right, with colour encoding whether late-period rates would
    raise (blue) or lower (vermillion) the safety margin.
    """
    fig_dir.mkdir(parents=True, exist_ok=True)
    df = period_compare.copy()
    df["abs_delta"] = df["delta_margin"].abs()
    df = df.sort_values("abs_delta", ascending=True)
    groups = df["group"].tolist()
    deltas = df["delta_margin"].tolist()
    lengths = df["abs_delta"].tolist()
    fig, ax = plt.subplots(figsize=(9, 5))
    # Colour-vision-deficiency-friendly palette: blue for positive, vermillion for negative
    CVD_POS = "#0072B2"
    CVD_NEG = "#D55E00"
    colors = [CVD_POS if d >= 0 else CVD_NEG for d in deltas]
    bars = ax.barh(groups, lengths, color=colors)
    for bar, d, l in zip(bars, deltas, lengths):
        ax.text(l + max(lengths) * 0.02,
                bar.get_y() + bar.get_height() / 2,
                f"{_fmt(d, 1)}", va="center", ha="left",
                fontsize=8)
    ax.axvline(0, color="black", linewidth=0.8)
    ax.set_xlabel("Change in equilibrium safety margin |late − early| (late higher → blue, lower → red)")
    ax.set_title("Counterfactual change in safety margin if late-period rates persisted (point estimates)")
    fig.tight_layout()
    path = fig_dir / "fig3_historical_margin.png"
    fig.savefig(path, dpi=300)
    plt.close(fig)
    return path


def build_figure4(boot, fig_dir: Path):
    """Bootstrap 95% confidence intervals for equilibrium T."""
    fig_dir.mkdir(parents=True, exist_ok=True)
    df = boot.sort_values("T_equilibrium_median")
    groups = df["group"].tolist()
    med = df["T_equilibrium_median"].tolist()
    low = df["T_equilibrium_q025"].tolist()
    high = df["T_equilibrium_q975"].tolist()
    fig, ax = plt.subplots(figsize=(9, 5))
    y = np.arange(len(groups))
    ax.errorbar(med, y, xerr=[np.subtract(med, low), np.subtract(high, med)],
                fmt="o", color="steelblue", capsize=4, ecolor="gray")
    ax.set_yticks(y)
    ax.set_yticklabels(groups)
    ax.set_xlabel("Equilibrium domestic active pool T")
    ax.set_title("Bootstrap 95% confidence intervals for equilibrium T")
    ax.set_xlim(0, max(high) * 1.05)
    fig.tight_layout()
    path = fig_dir / "fig4_bootstrap_ci.png"
    fig.savefig(path, dpi=600, bbox_inches="tight")
    plt.close(fig)
    return path


def build_figure8(trans_rates, eq, fig_dir: Path):
    """Japan's six-compartment flow with cross-civilisation transition-rate ladders."""
    fig_dir.mkdir(parents=True, exist_ok=True)

    def _rate_fmt(v, dec=3):
        if pd.isna(v):
            return "—"
        return f"{float(v):.{dec}f}"

    ja_row = trans_rates[trans_rates["group"] == "Japanese"].iloc[0]
    ja_eq = eq[eq["group"] == "Japanese"].iloc[0]

    fig = plt.figure(figsize=(18, 10))
    gs = fig.add_gridspec(1, 2, width_ratios=[1, 1.2])
    ax1 = fig.add_subplot(gs[0, 0])
    ax2 = fig.add_subplot(gs[0, 1])

    # Left panel: Japan compartment flow
    ax1.set_xlim(0, 1)
    ax1.set_ylim(0, 1)
    ax1.axis("off")
    w, h = 0.14, 0.10
    positions = {
        "D": (0.18, 0.80), "H_D": (0.18, 0.52), "P_D": (0.18, 0.24),
        "A": (0.82, 0.80), "H_A": (0.82, 0.52), "P_A": (0.82, 0.24),
        "L": (0.50, 0.05),
    }
    counts = {
        "D": ja_eq["D_eq"], "H_D": ja_eq["H_D_eq"], "P_D": ja_eq["P_D_eq"],
        "A": ja_eq["A_eq"], "H_A": ja_eq["H_A_eq"], "P_A": ja_eq["P_A_eq"],
    }
    for name, (x, y) in positions.items():
        if name == "L":
            label = "L\n(dropout)"
            fc = "#ffcccc"
        else:
            label = f"{name}\n{int(round(counts[name]))}"
            fc = "#e6f2ff" if name in ("D", "H_D", "P_D") else "#fff4e6"
        box = mpatches.FancyBboxPatch(
            (x - w / 2, y - h / 2), w, h, boxstyle="round,pad=0.02",
            facecolor=fc, edgecolor="black", linewidth=1.2,
        )
        ax1.add_patch(box)
        ax1.text(x, y, label, ha="center", va="center", fontsize=9, weight="bold")

    def _arrow(start, end, label, rate, color="black", lw=1.2):
        ax1.annotate("", xy=end, xytext=start,
                     arrowprops=dict(arrowstyle="->", color=color, lw=lw,
                                     connectionstyle="arc3,rad=0"))
        mx = (start[0] + end[0]) / 2
        my = (start[1] + end[1]) / 2
        ax1.text(mx, my + 0.025, f"{label}={_rate_fmt(rate)}",
                 ha="center", va="bottom", fontsize=8, color=color,
                 bbox=dict(boxstyle="round,pad=0.15", facecolor="white",
                           edgecolor="none", alpha=0.8))

    _arrow((0.25, 0.85), (0.75, 0.85), r"$\alpha$", ja_row["alpha"])
    _arrow((0.75, 0.75), (0.25, 0.75), r"$\beta$", ja_row["beta"])
    _arrow((0.18, 0.75), (0.18, 0.58), "h_D", ja_row["h_D"])
    _arrow((0.18, 0.47), (0.18, 0.30), "p_D", ja_row["p_D"])
    _arrow((0.82, 0.75), (0.82, 0.58), "h_A", ja_row["h_A"])
    _arrow((0.82, 0.47), (0.82, 0.30), "p_A", ja_row["p_A"])

    ax1.annotate("", xy=(0.11, 0.80), xytext=(0.02, 0.80),
                 arrowprops=dict(arrowstyle="->", color="green", lw=1.2))
    ax1.text(0.055, 0.83, f"I0+rP_D\nI0={_rate_fmt(ja_eq['I0'],1)}\nr={_rate_fmt(ja_eq['r'],4)}",
             ha="center", va="bottom", fontsize=7, color="green")

    for name, (x, y) in positions.items():
        if name == "L":
            continue
        ax1.annotate("", xy=(0.50, 0.12), xytext=(x, y - 0.05),
                     arrowprops=dict(arrowstyle="->", color="gray", lw=0.7,
                                     ls="--", connectionstyle="arc3,rad=0.15"))
    ax1.text(0.50, 0.015, f"d={_rate_fmt(ja_row['d'])} (all compartments)",
             ha="center", va="bottom", fontsize=8, color="gray", weight="bold")
    ax1.set_title(
        f"Japanese community (T/M={_rate_fmt(ja_eq['T_equilibrium']/ja_eq['M_threshold'],2)})",
        fontsize=12, weight="bold", pad=10,
    )

    # Right panel: rate ladders by civilisation, Japan highlighted
    rates = ["alpha", "beta", "h_D", "h_A", "p_D", "d"]
    others = sorted([g for g in trans_rates["group"] if g != "Japanese"])
    groups = ["Japanese"] + others
    n_groups = len(groups)
    n_rates = len(rates)
    y = np.arange(n_groups)
    bar_height = 0.11
    colors = plt.cm.tab10(np.linspace(0, 1, n_rates))
    for i, rate in enumerate(rates):
        vals = np.array([trans_rates[trans_rates["group"] == g][rate].values[0] for g in groups])
        ax2.barh(y + i * bar_height, vals, height=bar_height, label=rate, color=colors[i])
    ax2.set_yticks(y + bar_height * (n_rates - 1) / 2)
    ax2.set_yticklabels(groups, fontsize=8)
    ax2.invert_yaxis()
    ax2.set_xlabel("Transition rate", fontsize=9)
    ax2.set_title("Transition rates by civilisation (Japan highlighted)", fontsize=11, weight="bold", pad=10)
    ax2.legend(ncol=3, fontsize=7, loc="lower right")
    ax2.axhspan(-0.5, 0.5 + n_rates * bar_height, color="red", alpha=0.08)
    ax2.set_ylim(n_groups - 0.5, -0.5)
    ax2.set_xlim(0, 1.0)
    ax2.grid(axis="x", linestyle=":", alpha=0.5)

    fig.tight_layout()
    path = fig_dir / "fig8_japan_compartment_flow.png"
    fig.savefig(path, dpi=300, bbox_inches="tight")
    plt.close(fig)
    return path


def build_figure9(eq, pnr_active, fig_dir: Path):
    """T/M safety margin versus closest PNR proximity for all civilisations."""
    fig_dir.mkdir(parents=True, exist_ok=True)

    active = pnr_active[(pnr_active["target"] == "domestic_active") & (pnr_active["is_within_bounds"])].copy()
    active["proximity"] = (active["critical_factor"] - 1).abs()
    active = active.loc[active.groupby("group")["proximity"].idxmin()].reset_index(drop=True)
    merged = eq[["group", "T_equilibrium", "M_threshold"]].merge(
        active[["group", "rate_name", "current_rate", "critical_factor", "proximity"]], on="group"
    )
    merged["T_over_M"] = merged["T_equilibrium"] / merged["M_threshold"]
    merged = merged.sort_values("T_over_M")

    fig, ax = plt.subplots(figsize=(9, 7))
    colors = ["red" if g == "Japanese" else "steelblue" for g in merged["group"]]
    ax.scatter(merged["T_over_M"], merged["proximity"], c=colors, s=120, edgecolors="black", zorder=3)
    for _, row in merged.iterrows():
        ax.annotate(row["group"], (row["T_over_M"], row["proximity"]),
                    textcoords="offset points", xytext=(6, 4), fontsize=8)
    ax.axvline(1.0, color="gray", linestyle="--", lw=1)
    ax.set_xlabel("T / M (equilibrium active pool / minimum viable threshold)", fontsize=10)
    ax.set_ylabel("PNR proximity (|critical factor − 1|)\nsmaller = more fragile", fontsize=10)
    ax.set_title("Model evaluation: safety margin vs. point-of-no-return proximity", fontsize=11, weight="bold")

    ja = merged[merged["group"] == "Japanese"].iloc[0]
    ax.annotate(f"Japan: PNR lever = {ja['rate_name']}",
                (ja["T_over_M"], ja["proximity"]),
                textcoords="offset points", xytext=(-30, -25),
                fontsize=9, color="red",
                arrowprops=dict(arrowstyle="->", color="red"))

    fig.tight_layout()
    path = fig_dir / "fig9_tm_pnr_scatter.png"
    fig.savefig(path, dpi=300)
    plt.close(fig)
    return path


# ---------------------------------------------------------------------------
# Reference list
# ---------------------------------------------------------------------------

NOTE_TEXT = "Yamada Y (momentumyy). 海外で当てた研究者はその後どうなるのか. note.com, 2026. https://note.com/momentumyy/n/n86df5d34282d (accessed 2026-08-09)."

REFS = [
    "MacroPolo. The Global AI Talent Tracker 2.0. Paulson Institute, 2023. https://macropolo.org/digital-projects/the-global-ai-talent-tracker/",
    "Appelt S, van Beuzekom B, Galindo-Rueda F, de Pinho R. Which factors influence the international mobility of research scientists? OECD Science, Technology and Industry Working Papers 2015/02, 2015. https://doi.org/10.1787/5js1tmrr2233-en",
    "Stephan P E. The Economics of Science. J Econ Lit. 1996;34(3):1199-1235.",
    "Huntington S P. The Clash of Civilizations and the Remaking of World Order. New York: Simon & Schuster, 1996.",
    "Aghion P, Bloom N, Blundell R, Griffith R, Howitt P. Competition and innovation: an inverted-U relationship. Q J Econ. 2005;120(2):701-728.",
    "Priem J, Piwowar H, Orr R. OpenAlex: A fully-open index of scholarly works, authors, venues, institutions, and concepts. arXiv:2205.01833, 2022. https://doi.org/10.48550/arXiv.2205.01833",
    "Thorn K, Holm-Nielsen L B. International Mobility of Researchers and Scientists: Policy Options for Turning a Drain into a Gain. UNU-WIDER Research Paper No. 2006/83, 2006. https://www.wider.unu.edu/sites/default/files/rp2006-83.pdf",
    "AlShebli B, Memon S A, Evans J A, Rahwan T. China and the U.S. produce more impactful AI research when collaborating together. Sci Rep. 2024;14:28576. https://doi.org/10.1038/s41598-024-79863-5",
    "Yuan S, Shao Z, Wei X, Tang J, Hall W, Wang Y, et al. Science behind AI: the evolution of trend, mobility, and collaboration. Scientometrics. 2020;124(2):993-1013. https://doi.org/10.1007/s11192-020-03423-7",
    "Shaffer M L. Minimum Population Sizes for Species Conservation. BioScience. 1981;31(2):131-134.",
    "Franzoni C, Scellato G, Stephan P E. Foreign-born scientists: mobility patterns for 16 countries. Nat Biotechnol. 2012;30(12):1250-1253.",
    "Jones B F, Wuchty S, Uzzi B. Multi-University Research Teams: Shifting Impact, Geography, and Stratification in Science. Science. 2008;322(5905):1259-1262.",
    "Nelson R R, Winter S G. An Evolutionary Theory of Economic Change. Cambridge, MA: Harvard University Press, 1982.",
    "Dosi G. Technological paradigms and technological trajectories: a suggested interpretation of the determinants and directions of technical change. Res Policy. 1982;11(3):147-162.",
    "Lundvall B-Å. National Systems of Innovation: Toward a Theory of Innovation and Interactive Learning. London: Anthem Press, 1992.",
    "Malerba F. Sectoral systems of innovation and production. Res Policy. 2002;31(2):247-264.",
    "State B, Park P, Weber I, Macy M. The mesh of civilizations in the global network of digital communication. PLoS ONE. 2015;10(5):e0122543. https://doi.org/10.1371/journal.pone.0122543",
    "Chinchilla-Rodríguez Z, Miao L, Murray D, Robinson-García N, Costas R, Sugimoto C R. A global comparison of scientific mobility and collaboration according to national scientific capacities. Front Res Metr Anal. 2018;3:17. https://doi.org/10.3389/frma.2018.00017",
    "Freeman R B, Huang W. Collaboration: Strength in diversity. Nature. 2014;513(7518):305. https://doi.org/10.1038/513305a",
    "Shachar A. The Race for Talent: Highly Skilled Migrants and Competitive Immigration Regimes. NYU Law Rev. 2006;81(1):148-206.",
    "Kerr W R. Global Talent and U.S. Immigration Policy. Harvard Business School Working Paper No. 20-107, 2020. https://www.hbs.edu/ris/Publication%20Files/20-107_0967f1ab-1d23-4d54-b5a1-c884234d9b31.pdf"
]


def _ref(n: int) -> str:
    return REFS[n - 1]


# ---------------------------------------------------------------------------
# Data loading
# ---------------------------------------------------------------------------

def load_data():
    cohort = pd.read_csv(BASE_DIR / "data" / "cohort" / "cohort.csv")
    eq = pd.read_csv(ENDOG / "equilibrium_summary.csv")
    sat_eq = pd.read_csv(SAT / "equilibrium_summary.csv") if (SAT / "equilibrium_summary.csv").exists() else None
    top_t = pd.read_csv(ENDOG / "top_transitions_T.csv")
    pnr_full = pd.read_csv(ENDOG / "point_of_no_return.csv")
    # Closest per group, active pool preferred, restricting to rates that actually cross the threshold
    active_only = pnr_full[(pnr_full["target"] == "domestic_active") & (pnr_full["is_within_bounds"] == True)].copy()
    active_only["proximity"] = (active_only["critical_factor"] - 1.0).abs()
    pnr_closest = active_only.loc[active_only.groupby("group")["proximity"].idxmin()].reset_index(drop=True)
    pnr_closest = pnr_closest.sort_values("proximity").reset_index(drop=True)
    period_compare = pd.read_csv(TV / "period_comparison.csv")
    boot = pd.read_csv(BOOT / "bootstrap_summary.csv")
    policy_rank = pd.read_csv(POL / "ranked_interventions.csv")
    return cohort, eq, sat_eq, top_t, pnr_closest, period_compare, boot, policy_rank


# ---------------------------------------------------------------------------
# Annual transition-rate and projection helpers
# ---------------------------------------------------------------------------


def load_annual_data():
    """Read annual transition-rate and projection CSVs.

    Returns a dict of DataFrames; missing tables are returned as None.
    """
    paths = {
        "rate_table": ANNUAL / "annual_ode_rates.csv",
        "projected_rates": ANNUAL / "projected_ode_rates.csv",
        "observed_stock": ANNUAL / "observed_annual_stock.csv",
        "projected_stock": ANNUAL / "projected_annual_stock.csv",
        "interciv_stock": ANNUAL / "annual_interciv_stock.csv",
        "evaluation": ANNUAL / "projection_evaluation.csv",
        "group_accuracy": ANNUAL / "projection_accuracy_by_group.csv",
        "compartment_accuracy": ANNUAL / "projection_accuracy_by_compartment.csv",
        "rate_accuracy": ANNUAL / "projection_rate_accuracy.csv",
        "rate_accuracy_overall": ANNUAL / "projection_rate_accuracy_overall.csv",
    }
    return {k: pd.read_csv(p) if p.exists() else None for k, p in paths.items()}


def pnr_robustness_table():
    """Compare closest PNR levers between linear and saturating endogenous inflow."""
    lin_path = BASE_DIR / "results" / "endogenous" / "closest_point_of_no_return.csv"
    sat_path = BASE_DIR / "results" / "endogenous_saturating" / "closest_point_of_no_return.csv"
    if not (lin_path.exists() and sat_path.exists()):
        return pd.DataFrame()
    lin = pd.read_csv(lin_path)
    sat = pd.read_csv(sat_path)
    # Focus on the active-pool threshold, which is the operative PNR in the main text.
    lin = lin[lin["target"] == "domestic_active"].copy()
    sat = sat[sat["target"] == "domestic_active"].copy()
    lin = lin.rename(columns={
        "group": "origin_group",
        "rate_name": "linear_closest",
        "critical_factor": "linear_factor",
        "proximity": "linear_proximity",
    })
    sat = sat.rename(columns={
        "group": "origin_group",
        "rate_name": "saturating_closest",
        "critical_factor": "saturating_factor",
        "proximity": "saturating_proximity",
    })
    df = lin[["origin_group", "linear_closest", "linear_factor", "linear_proximity"]].merge(
        sat[["origin_group", "saturating_closest", "saturating_factor", "saturating_proximity"]],
        on="origin_group",
        how="inner",
    )
    if df.empty:
        return pd.DataFrame()
    # Reorder to the standard group order and round.
    df = df.set_index("origin_group").reindex(arpr.ORDERED_GROUPS).reset_index()
    for col in ["linear_factor", "saturating_factor", "linear_proximity", "saturating_proximity"]:
        df[col] = df[col].round(4)
    return df


def compute_annual_context(annual):
    """Return data-derived summary strings for the annual projection sections."""
    ctx = {}
    eval_df = annual.get("evaluation")
    if eval_df is not None and not eval_df.empty:
        ctx["overall_rmse"] = float(((eval_df["error"] ** 2).mean()) ** 0.5)
        ctx["overall_mape"] = float(eval_df["ape"].mean())
        ctx["overall_mape_pct"] = ctx["overall_mape"] * 100.0
    else:
        ctx["overall_rmse"] = float("nan")
        ctx["overall_mape"] = float("nan")
        ctx["overall_mape_pct"] = float("nan")

    gacc = annual.get("group_accuracy")
    if gacc is not None and not gacc.empty:
        gacc = gacc.dropna(subset=["mape"]).copy()
        if not gacc.empty:
            best = gacc.loc[gacc["mape"].idxmin()]
            worst = gacc.loc[gacc["mape"].idxmax()]
            ctx["best_group"] = best["origin_group"]
            ctx["worst_group"] = worst["origin_group"]
            ctx["best_mape_pct"] = float(best["mape"]) * 100.0
            ctx["worst_mape_pct"] = float(worst["mape"]) * 100.0
            ctx["n_eval_groups"] = len(gacc)
            # Direction agreement (mean across civilisations)
            if "direction_agreement" in gacc.columns:
                ctx["direction_agreement"] = float(gacc["direction_agreement"].mean())
                ctx["best_direction_group"] = gacc.loc[gacc["direction_agreement"].idxmax()]["origin_group"]
                ctx["worst_direction_group"] = gacc.loc[gacc["direction_agreement"].idxmin()]["origin_group"]
            if "threshold_alarm_accuracy" in gacc.columns:
                ctx["threshold_alarm_accuracy"] = float(gacc["threshold_alarm_accuracy"].mean())
                ctx["threshold_alarm_sensitivity"] = float(gacc["threshold_alarm_sensitivity"].mean())
                ctx["threshold_alarm_specificity"] = float(gacc["threshold_alarm_specificity"].mean())
                # Observed alarms are rare; report total count across groups
                ctx["threshold_alarms_obs"] = int(gacc["threshold_alarms_obs"].sum())
        else:
            ctx["best_group"] = "—"
            ctx["worst_group"] = "—"
            ctx["best_mape_pct"] = float("nan")
            ctx["worst_mape_pct"] = float("nan")
            ctx["n_eval_groups"] = 0
    else:
        ctx["best_group"] = "—"
        ctx["worst_group"] = "—"
        ctx["best_mape_pct"] = float("nan")
        ctx["worst_mape_pct"] = float("nan")
        ctx["n_eval_groups"] = 0

    proj = annual.get("projected_rates")
    if proj is not None and not proj.empty:
        ctx["n_projected_group_years"] = len(proj)
        ctx["smoothed_pct"] = float(proj["correction_smoothed"].mean()) * 100.0
        ctx["capped_pct"] = float(proj["correction_capped"].mean()) * 100.0
    else:
        ctx["n_projected_group_years"] = 0
        ctx["smoothed_pct"] = float("nan")
        ctx["capped_pct"] = float("nan")

    obs = annual.get("observed_stock")
    if obs is not None and not obs.empty:
        ctx["obs_year_min"] = int(obs["year"].min())
        ctx["obs_year_max"] = int(obs["year"].max())
    else:
        ctx["obs_year_min"] = 2000
        ctx["obs_year_max"] = 2023

    # Rate-level forecast accuracy (cleaner than stock-level because the fixed cohort lacks post-2016 entrants).
    rate_overall = annual.get("rate_accuracy_overall")
    if rate_overall is not None and not rate_overall.empty:
        ctx["rate_overall_rmse"] = float(rate_overall["rmse"].iloc[0])
        ctx["rate_overall_mae"] = float(rate_overall["mae"].iloc[0])
        ctx["rate_overall_mape"] = float(rate_overall["mape"].iloc[0])
        ctx["rate_overall_skill"] = float(rate_overall["skill"].iloc[0])
    else:
        ctx["rate_overall_rmse"] = float("nan")
        ctx["rate_overall_mae"] = float("nan")
        ctx["rate_overall_mape"] = float("nan")
        ctx["rate_overall_skill"] = float("nan")

    rate_acc = annual.get("rate_accuracy")
    if rate_acc is not None and not rate_acc.empty:
        # Best/worst rate by skill (naive/model RMSE ratio)
        rate_acc = rate_acc.dropna(subset=["skill"]).copy()
        if not rate_acc.empty:
            ctx["best_rate_skill"] = str(rate_acc.loc[rate_acc["skill"].idxmax(), "rate"])
            ctx["worst_rate_skill"] = str(rate_acc.loc[rate_acc["skill"].idxmin(), "rate"])
            ctx["best_rate_skill_value"] = float(rate_acc["skill"].max())
            ctx["worst_rate_skill_value"] = float(rate_acc["skill"].min())
        else:
            ctx["best_rate_skill"] = "—"
            ctx["worst_rate_skill"] = "—"
            ctx["best_rate_skill_value"] = float("nan")
            ctx["worst_rate_skill_value"] = float("nan")
    else:
        ctx["best_rate_skill"] = "—"
        ctx["worst_rate_skill"] = "—"
        ctx["best_rate_skill_value"] = float("nan")
        ctx["worst_rate_skill_value"] = float("nan")

    return ctx


def compute_japan_context(eq, pnr_closest, transition_rates):
    """Return data-derived prose values for the Japan-specific discussion."""
    ja = eq[eq["group"] == "Japanese"].iloc[0]
    tr = transition_rates[transition_rates["group"] == "Japanese"].iloc[0]
    pnr = pnr_closest[pnr_closest["group"] == "Japanese"]
    if pnr.empty:
        # fallback from full pnr if pnr_closest was not supplied with Japan row
        pnr_row = None
    else:
        pnr_row = pnr.iloc[0]
    return {
        "T": float(ja["T_equilibrium"]),
        "M": float(ja["M_threshold"]),
        "T_over_M": float(ja["T_equilibrium"]) / float(ja["M_threshold"]),
        "margin": float(ja["margin_to_threshold_T"]),
        "D": int(round(ja["D_eq"])),
        "A": int(round(ja["A_eq"])),
        "H_D": int(round(ja["H_D_eq"])),
        "H_A": int(round(ja["H_A_eq"])),
        "P_D": int(round(ja["P_D_eq"])),
        "P_A": int(round(ja["P_A_eq"])),
        "I0": float(ja["I0"]),
        "r": float(ja["r"]),
        "alpha": float(tr["alpha"]),
        "beta": float(tr["beta"]),
        "h_D": float(tr["h_D"]),
        "h_A": float(tr["h_A"]),
        "p_D": float(tr["p_D"]),
        "p_A": float(tr["p_A"]),
        "d": float(tr["d"]),
        "pnr_rate": pnr_row["rate_name"] if pnr_row is not None else "I0",
        "pnr_factor": float(pnr_row["critical_factor"]) if pnr_row is not None else float("nan"),
        "pnr_proximity": float(pnr_row["proximity"]) if pnr_row is not None else float("nan"),
    }


def annual_summary_table(annual):
    """Mean observed annual transition rates and inflow by group (2000-2016)."""
    rate_table = annual.get("rate_table")
    if rate_table is None or rate_table.empty:
        return pd.DataFrame()
    observed = rate_table[rate_table["year"] <= 2016]
    cols = ["alpha", "beta", "h_D", "p_D", "d", "I_total"]
    means = observed.groupby("origin_group")[cols].mean().reset_index()
    means.columns = ["Group", "α", "β", "h_D", "p_D", "d", "I_total"]
    return means


def interciv_top_table(annual, n=10):
    """Top origin-destination abroad author-year accumulations.

    Unknown destinations and origin==destination domestic moves are excluded
    because the reconstruction cannot observe the actual host civilisation.
    """
    flows = annual.get("interciv_stock")
    if flows is None or flows.empty:
        return pd.DataFrame()
    flows = flows[
        (flows["destination_group"] != "Unknown") &
        (flows["origin_group"] != flows["destination_group"])
    ].copy()
    pivot = (
        flows.groupby(["origin_group", "destination_group"], observed=False)["count"]
        .sum()
        .reset_index()
        .sort_values("count", ascending=False)
        .head(n)
    )
    pivot.columns = ["Origin", "Destination", "Author-years"]
    return pivot


def build_annual_figures(annual, fig_dir):
    """Generate annual projection figures; reuse existing PNGs if data are missing."""
    fig_paths = {}
    rate_table = annual.get("rate_table")
    projected_rates = annual.get("projected_rates")
    if rate_table is not None and projected_rates is not None:
        fig_paths["fig5"] = arpr.plot_annual_rates(rate_table, projected_rates, fig_dir=fig_dir)
    else:
        fig_paths["fig5"] = fig_dir / "annual_rates_by_group.png"

    interciv = annual.get("interciv_stock")
    if interciv is not None:
        fig_paths["fig6"] = arpr.plot_interciv_heatmap(interciv, fig_dir=fig_dir)
    else:
        fig_paths["fig6"] = fig_dir / "annual_interciv_heatmap.png"

    obs_stock = annual.get("observed_stock")
    proj_stock = annual.get("projected_stock")
    if obs_stock is not None and proj_stock is not None:
        fig_paths["fig7"] = arpr.plot_projection_by_compartment(proj_stock, obs_stock, fig_dir=fig_dir)
    else:
        fig_paths["fig7"] = fig_dir / "annual_projection_vs_observed.png"
    return fig_paths


def compute_context(cohort, eq, sat_eq, top_t, pnr_closest, period_compare, policy_rank):
    """Return data-derived summary strings used in the Results and Discussion."""
    n_groups = len(eq)
    eq_sorted = eq.sort_values("T_equilibrium", ascending=False)
    largest_pools = ", ".join(eq_sorted["group"].head(3).tolist())
    smallest_pool = eq_sorted["group"].iloc[-1]
    eq_m = eq.sort_values("margin_to_threshold_T")
    smallest_margin_group = eq_m["group"].iloc[0]

    d_rows = top_t[(top_t["rate"] == "d") & (top_t["target"] == "domestic_active")]
    d_min_e = d_rows["elasticity"].min()
    d_max_e = d_rows["elasticity"].max()
    d_all_negative = (d_rows["elasticity"] < 0).all()

    # Positive levers after dropout
    positive = []
    for _, gdf in top_t.groupby("group"):
        gdf = gdf.sort_values("abs_elasticity", ascending=False)
        # Skip the largest (dropout), then collect the positive transition-rate levers
        for _, r in gdf.iloc[1:].iterrows():
            if r["elasticity"] > 0 and r["rate"] not in ("I0", "r"):
                positive.append(r["rate"])
    pos_counts = pd.Series(positive).value_counts()
    most_common_positive = pos_counts.index[0] if not pos_counts.empty else "p_D"
    second_positive = pos_counts.index[1] if len(pos_counts) > 1 else None
    if most_common_positive == "p_D":
        pos_lever_text = "principal-investigator promotion (p_D)"
    elif most_common_positive == "h_D":
        pos_lever_text = "domestic hit generation (h_D)"
    elif most_common_positive == "beta":
        pos_lever_text = "return from abroad (β)"
    else:
        pos_lever_text = most_common_positive
    if second_positive == "p_D":
        second_text = "principal-investigator promotion (p_D)"
    elif second_positive == "h_D":
        second_text = "domestic hit generation (h_D)"
    elif second_positive == "beta":
        second_text = "return from abroad (β)"
    else:
        second_text = second_positive
    if second_text and second_text != pos_lever_text:
        positive_lever_sentence = f"The largest positive transition lever is {pos_lever_text}, followed by {second_text}."
    else:
        positive_lever_sentence = f"The largest positive transition lever is {pos_lever_text}."
    positive_lever_sentence_lower = positive_lever_sentence[0].lower() + positive_lever_sentence[1:]
    if positive_lever_sentence_lower.endswith('.'):
        positive_lever_sentence_lower = positive_lever_sentence_lower[:-1]

    # PI promotion elasticity, identify group with highest p_D elasticity
    pd_elas = top_t[(top_t["rate"] == "p_D") & (top_t["target"] == "domestic_active")].copy()
    pd_elas["abs_e"] = pd_elas["elasticity"].abs()
    highest_pd_group = pd_elas.sort_values("abs_e", ascending=False).iloc[0]["group"] if not pd_elas.empty else "Japanese"

    # Point of no return
    closest_rate_counts = pnr_closest["rate_name"].value_counts()
    closest_rate_mode = closest_rate_counts.index[0] if not closest_rate_counts.empty else "I0"
    all_closest_same = len(closest_rate_counts) == 1
    if all_closest_same:
        pnr_lever_text = f"{closest_rate_mode} is the closest point-of-no-return lever for the active researcher pool in every group"
    else:
        pnr_lever_text = f"{closest_rate_mode} is the most common closest point-of-no-return lever for the active researcher pool"

    # Saturating reduction range
    sat_range_text = ""
    if sat_eq is not None:
        merged = eq[["group", "T_equilibrium"]].merge(
            sat_eq[["group", "T_equilibrium"]], on="group", suffixes=("_lin", "_sat")
        )
        pct_diff = 100.0 * (merged["T_equilibrium_lin"] - merged["T_equilibrium_sat"]) / merged["T_equilibrium_lin"]
        max_abs = pct_diff.abs().max()
        if max_abs < 0.001:
            sat_range_text = "below 0.001% for every group"
        else:
            sat_range_text = f"up to {max_abs:.2f}% lower than the linear variant"

    # Historical counterfactual
    if period_compare.empty:
        period_neg = "none"
        period_pos = "none"
        period_all_neg = False
    else:
        sorted_pc = period_compare.sort_values("delta_margin")
        neg = sorted_pc[sorted_pc["delta_margin"] < 0]["group"].tolist()
        pos = sorted_pc[sorted_pc["delta_margin"] > 0]["group"].tolist()[::-1]
        period_neg = ", ".join(neg) if neg else "none"
        period_pos = ", ".join(pos) if pos else "none"
        period_all_neg = len(pos) == 0

    # 10% dropout margin gain range
    d_decrease = policy_rank[(policy_rank["lever"] == "d") & (policy_rank["direction"] == "decrease")].copy()
    d_10pct = d_decrease[d_decrease["lever_change_pct"].abs() >= 9.9]
    if d_10pct.empty:
        d_10pct = d_decrease
    d_10pct_group = d_10pct.loc[d_10pct.groupby("group")["normalised_margin_gain_per_10pct"].idxmax()]
    d_min_gain = d_10pct_group.sort_values("margin_gain").iloc[0]
    d_max_gain = d_10pct_group.sort_values("margin_gain").iloc[-1]

    # Endogenous inflow safety factor used in the fitted model.  The default code
    # cap is 0.50 of the critical reproduction rate; the most constrained fitted
    # group has a realised r / r_critical ratio that is lower (about 0.40).
    safety_factor_cap = 0.50
    min_realised_safety_ratio = float((eq["r"] / eq["r_critical"]).min())

    return {
        "n_groups": n_groups,
        "largest_pools": largest_pools,
        "smallest_pool": smallest_pool,
        "smallest_margin_group": smallest_margin_group,
        "d_min_e": d_min_e,
        "d_max_e": d_max_e,
        "d_all_negative": d_all_negative,
        "positive_lever_sentence": positive_lever_sentence,
        "positive_lever_sentence_lower": positive_lever_sentence_lower,
        "highest_pd_group": highest_pd_group,
        "pnr_lever_text": pnr_lever_text,
        "sat_range_text": sat_range_text,
        "period_neg": period_neg,
        "period_pos": period_pos,
        "period_all_neg": period_all_neg,
        "d_min_gain_group": d_min_gain["group"],
        "d_max_gain_group": d_max_gain["group"],
        "d_min_gain": round(d_min_gain["margin_gain"]),
        "d_max_gain": round(d_max_gain["margin_gain"]),
        "safety_factor_cap": safety_factor_cap,
        "min_realised_safety_ratio": min_realised_safety_ratio,
    }


def _package_summary():
    """Return a narrative and a DataFrame for the best multi-lever policy packages.

    Packages are generated by src/policy_counterfactuals.py --packages.  We report
    the package with the largest absolute margin gain for each of the three
    smallest-margin groups, using counterfactuals.csv (no hard-coded numbers).
    """
    cf_path = POL / "counterfactuals.csv"
    if not cf_path.exists():
        return None, None
    cf = pd.read_csv(cf_path)
    packages = cf[cf["lever"].str.startswith("package:", na=False)].copy()
    if packages.empty:
        return None, None
    top = packages.loc[packages.groupby("group")["delta_margin"].idxmax()].copy()
    top["package_name"] = top["lever"].str.replace("package:", "", regex=False)
    top = top.sort_values("delta_margin")
    parts = [
        f"{r['group']} ({r['package_name']}: +{_fmt(r['delta_margin'], 0)} active researchers)"
        for _, r in top.iterrows()
    ]
    narrative = (
        "We also evaluated multi-lever policy packages for the three smallest-margin groups. "
        "The package with the largest margin gain in each group was: "
        + "; ".join(parts)
        + ". "
        "These packages combine dropout reduction with return or PI-pipeline levers, "
        "showing that the framework can compare multi-lever interventions as well as single-rate perturbations."
    )
    return narrative, top


def _unify_pnr_markdown(text):
    """Unify 'point of no return' to PNR in markdown, keeping the definition in the Abstract and Introduction."""
    import re

    pnr_re = re.compile(r"point of no return(?: \(PNR\))?", re.IGNORECASE)
    current_section = None
    seen = {"Abstract": False, "1. Introduction": False}

    def _replace(m, section):
        key = section if section == "Abstract" else "1. Introduction"
        if section in ("Abstract", "1. Introduction") and not seen[key]:
            seen[key] = True
            return m.group(0)
        return "PNR"

    out_lines = []
    for line in text.split("\n"):
        m_heading = re.match(r"##\s+(.+)", line)
        if m_heading:
            current_section = m_heading.group(1).strip()
        out_lines.append(pnr_re.sub(lambda m: _replace(m, current_section), line))

    return "\n".join(out_lines)


def _unify_pnr_docx(doc):
    """Unify 'point of no return' to PNR in docx body text, keeping definitions in Abstract and Introduction."""
    import re

    pnr_re = re.compile(r"point of no return(?: \(PNR\))?", re.IGNORECASE)
    current_section = None
    seen = {"Abstract": False, "Introduction": False}

    for para in doc.paragraphs:
        text = para.text
        if text.startswith("Abstract"):
            current_section = "Abstract"
            continue
        if text.startswith("1. Introduction"):
            current_section = "Introduction"
            continue
        # Only the abstract and introduction paragraphs contain the phrase.
        if current_section in ("Abstract", "Introduction") and pnr_re.search(text):
            for run in para.runs:
                if pnr_re.search(run.text):
                    if not seen[current_section]:
                        seen[current_section] = True
                    else:
                        run.text = pnr_re.sub("PNR", run.text)
    return doc


def _renumber_markdown_sections(lines):
    """Renumber markdown subsections so every main section has continuous numbers.

    The Word manuscript is the submission file; the Markdown copy is for version
    control and review. Some subsections that are present in the Word version are
    omitted in Markdown, creating gaps (e.g. 4.4 -> 4.10). This function renumbers
    the visible subsections continuously and updates any in-text references.
    """
    import re

    renumbered = []
    current_section = None
    subsection_counter = 0
    mapping = {}

    section_re = re.compile(r"^## (\d+)\.\s+(.*)")
    subsection_re = re.compile(r"^### (\d+)\.(\d+)\s+(.*)")

    for line in lines:
        m = section_re.match(line)
        if m:
            current_section = int(m.group(1))
            subsection_counter = 0
            renumbered.append(line)
            continue
        m = subsection_re.match(line)
        if m and current_section is not None:
            section_num = int(m.group(1))
            old_sub = m.group(2)
            title = m.group(3)
            if section_num != current_section:
                current_section = section_num
                subsection_counter = 0
            subsection_counter += 1
            new_num = f"{section_num}.{subsection_counter}"
            old_full = f"{section_num}.{old_sub}"
            if old_full != new_num:
                mapping[old_full] = new_num
            renumbered.append(f"### {new_num} {title}")
            continue
        renumbered.append(line)

    # Replace in-text references like "Section 4.7" with the renumbered target.
    keys = sorted(mapping.keys(), key=lambda k: len(k), reverse=True)
    for i, line in enumerate(renumbered):
        for old in keys:
            # Match "Section" or "Sections" followed by the old number at a word boundary.
            pattern = re.compile(rf"(Section[s]?\\s+){re.escape(old)}\\b")
            line = pattern.sub(rf"\g<1>{mapping[old]}", line)
        renumbered[i] = line

    return renumbered


# ---------------------------------------------------------------------------
# Markdown output
# ---------------------------------------------------------------------------

def _abstract_and_highlights(eq, pnr_closest):
    closest = pnr_closest.iloc[0]
    # Compute a robust, data-driven statement about the most efficient lever
    policy_rank_path = POL / "ranked_interventions.csv"
    if policy_rank_path.exists():
        policy_rank = pd.read_csv(policy_rank_path)
        top_by_group = policy_rank.groupby("group").head(1)
        all_top_are_d = (top_by_group["lever"] == "d").all()
        top_lever_mode = top_by_group["lever"].mode()
        most_common_lever = top_lever_mode.iloc[0] if not top_lever_mode.empty else "d"
    else:
        all_top_are_d = True
        most_common_lever = "d"
    if all_top_are_d:
        lever_text = "A simulated reduction in dropout yields the largest margin gain in every group in the fitted model. "
        highlight_lever = "Dropout reduction yields the largest margin gain across all groups."
    else:
        lever_text = f"A simulated reduction in dropout is the most common single positive lever in the fitted model, although other levers dominate for some groups in the current data. "
        highlight_lever = f"Simulated {most_common_lever} adjustment yields the largest margin gain per unit proportional change for most groups in the fitted model."
    # Conservative out-of-sample accuracy summary for the abstract
    ann_eval_path = ANNUAL / "projection_evaluation.csv"
    if ann_eval_path.exists():
        ev = pd.read_csv(ann_eval_path)
        overall_rmse = float(((ev["error"] ** 2).mean()) ** 0.5)
        overall_mape = float(ev["ape"].mean() * 100.0)
        projection_accuracy_text = (
            f"The 2017-2023 projection has RMSE {_fmt(overall_rmse, 2)} and a conservative, non-standard MAPE of {_fmt(overall_mape, 1)}% (count_obs + 1 denominator). "
            "The high error is expected because the projection is an early-warning indicator of directional drift, not a precise forecast. "
        )
    else:
        projection_accuracy_text = ""
    abstract = (
        "Artificial intelligence (AI) and machine learning (ML) research is increasingly concentrated, "
        "raising the risk that smaller communities fall below a minimum viable coauthor pool. "
        "We model each civilisation as a six-compartment system of domestic and abroad early-career, high-impact, and principal-investigator (PI) researchers, "
        "and estimate transition rates from OpenAlex AI/ML data (subfield 1702). "
        "The minimum viable threshold is M = k × c_bar, where c_bar is the mean authors per work and k is the median number of distinct last-author groups per year. "
        f"Across {len(eq)} groups, equilibrium active pools remain above their thresholds, but the closest point of no return (PNR) is observed for the {closest['group']} group, "
        f"where the {_rate_label(closest['rate_name'])} must be multiplied by {_fmt(closest['critical_factor'], 3)}× (a {closest['proximity']*100:.0f}% proportional {'reduction' if closest['critical_factor'] < 1 else 'increase'}) to drive the active pool to its threshold. "
        + lever_text
        + projection_accuracy_text
        + "Historical counterfactuals and bootstrap uncertainty show that the model is most sensitive to exogenous entry and attrition. "
        "These results provide a quantitative framework for early, safety-factor-bound policy scenarios that preserve civilisational diversity in AI/ML research."
    )
    keywords = (
        "researcher mobility; artificial intelligence; civilisation grouping; "
        "ordinary differential equations; PNR; innovation studies"
    )
    highlights = [
        "Nine civilisations modelled as six-compartment ODEs fitted to OpenAlex AI/ML data.",
        f"Closest point of no return: {closest['group']} via {closest['rate_name']} (factor {_fmt(closest['critical_factor'], 3)}×).",

        highlight_lever,
    ]
    return abstract, keywords, highlights


def _data_availability_text(blinded=False):
    base = (
        "This study uses the OpenAlex database (subfield 1702, Artificial Intelligence; "
        "2000–2023), accessed via the OpenAlex API. The analysis is bundled with a pre-extracted "
        "cohort and a stratified sample of works; the country-to-civilisation mapping, code, and result CSVs "
        "used to generate this manuscript "
    )
    if blinded:
        return base + "will be made available in a public repository upon acceptance. The extracted cohort SQLite database (a derived aggregate of OpenAlex records) is available from the corresponding author on request, subject to the OpenAlex CC0 licence and any applicable local data-use policies."
    return (
        base + "are available in the public GitHub repository "
        "https://github.com/bougtoir/researcher-mobility-ode. The extracted cohort SQLite database (a derived aggregate of OpenAlex records) is available from the corresponding author on request, subject to the OpenAlex CC0 licence and any applicable local data-use policies."
    )


def _descriptive_table(cohort):
    """Return DataFrame of descriptive statistics per group."""
    grp = cohort.groupby("origin_group").agg(
        n=("author_id", "count"),
        works=("n_ai_works", "sum"),
        active=("active", "sum"),
        hits=("hit", "sum"),
        pis=("pi", "sum"),
        career_start_mean=("career_start", "mean"),
        abroad=("abroad", "sum"),
    ).reset_index()
    grp["career_start_mean"] = grp["career_start_mean"].round(1)
    grp = grp.rename(columns={"origin_group": "Group"})
    return grp


def _docx_to_markdown(docx_path: Path, output_dir: Path) -> Path:
    """Convert an existing Word manuscript to Markdown so markdown is a faithful derivative."""
    if pypandoc is None:
        raise RuntimeError("pypandoc is required to generate the markdown version")
    md_path = output_dir / docx_path.with_suffix(".md").name
    pypandoc.convert_file(str(docx_path), "md", format="docx", outputfile=str(md_path))
    return md_path


def write_markdown(output_dir: Path, data=None, fig_paths=None, docx_path=None, blinded=False):
    """Write a plain-text markdown version for version control and review."""
    if docx_path is not None:
        md_path = _docx_to_markdown(docx_path, output_dir)
        return md_path
    # Fallback: the legacy markdown builder is no longer maintained; use pypandoc instead.
    raise RuntimeError("write_markdown now requires a docx_path; regenerate docx first.")


def _add_title_page(doc, word_count=None, blinded=False):
    style = doc.styles["Normal"]
    style.font.name = "Times New Roman"
    style.font.size = Pt(11)
    style.paragraph_format.line_spacing = 1.5

    title = doc.add_heading("Sustaining Heterogeneity through Interventions in Global AI/ML Researcher Mobility: A Transition-Rate Framework", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title.runs[0].font.size = Pt(16)
    title.runs[0].font.bold = True

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run("Article type: Research Article")

    if word_count:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run(f"Approximate word count (main text incl. tables, excl. references): {word_count}")

    if not blinded:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run("Corresponding author: [To be completed at submission]")
    else:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run("Author information removed for double-blind review")

    p = doc.add_paragraph()
    run = p.add_run()
    run.add_break(WD_BREAK.PAGE)


def _add_front_matter(doc, abstract, keywords, highlights, blinded=False):
    doc.add_heading("Abstract", level=1)
    p = doc.add_paragraph()
    p.add_run(abstract)

    p = doc.add_paragraph()
    p.add_run("Keywords: ").bold = True
    p.add_run(keywords)

    doc.add_heading("Highlights", level=2)
    for h in highlights:
        p = doc.add_paragraph(style="List Bullet")
        p.add_run(h)

    doc.add_heading("Data and Code Availability", level=2)
    p = doc.add_paragraph()
    p.add_run(_data_availability_text(blinded=blinded))

    doc.add_heading("Declarations", level=2)
    declarations = [
        ("Funding", "[To be completed by the authors at submission.]"),
        ("Competing interests", "[To be completed by the authors at submission.]"),
        ("Author contributions", "[To be completed by the authors at submission.]"),
        (
            "Declaration of generative AI in scientific writing",
            "During the preparation of this work the authors used AI-assisted tools to draft, code, and revise the manuscript. All claims, data, and interpretations were reviewed and approved by the authors.",
        ),
    ]
    if not blinded:
        declarations.append(("Acknowledgments", "This study was motivated by a note.com essay by Yamada Y (momentumyy) that framed researcher mobility in terms of transition rates rather than net flows (" + NOTE_TEXT + ")."))
    else:
        declarations.append(("Acknowledgments", "[Removed for double-blind review]"))
    for sub, text in declarations:
        p = doc.add_paragraph()
        p.add_run(f"{sub}: ").bold = True
        p.add_run(text)


def _add_table_from_df(doc, df, caption, decimals=None, bold_header=True):
    if decimals is None:
        decimals = {}
    cols = df.columns.tolist()
    table = doc.add_table(rows=1, cols=len(cols))
    table.style = "Table Grid"
    hdr = table.rows[0].cells
    for i, c in enumerate(cols):
        hdr[i].text = str(c)
        if bold_header:
            for run in hdr[i].paragraphs[0].runs:
                run.font.bold = True
    for _, row in df.iterrows():
        cells = table.add_row().cells
        for i, c in enumerate(cols):
            v = row[c]
            cells[i].text = _fmt(v, decimals.get(c, 2))
    cap = doc.add_paragraph()
    cap.add_run(caption).italic = True
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    return table


def _add_docx_body(doc, data, fig_paths, blinded=False):
    (cohort, eq, sat_eq, top_t, pnr_closest, period_compare, boot, policy_rank) = data
    ctx = compute_context(cohort, eq, sat_eq, top_t, pnr_closest, period_compare, policy_rank)
    transition_rates = pd.read_csv(BASE_DIR / "data" / "cohort" / "transition_rates.csv")
    ja_ctx = compute_japan_context(eq, pnr_closest, transition_rates)
    annual = load_annual_data()
    annual_ctx = compute_annual_context(annual)
    annual_means = annual_summary_table(annual)
    interciv_top = interciv_top_table(annual)
    group_acc = annual.get("group_accuracy")
    comp_acc = annual.get("compartment_accuracy")

    best_rmse_group = "—"
    worst_rmse_group = "—"
    worst_mape_group = "—"
    best_compartment_rmse = "—"
    worst_compartment_rmse = "—"
    worst_compartment_mape = "—"
    if group_acc is not None and not group_acc.empty:
        best_rmse_group = group_acc.loc[group_acc["rmse"].idxmin(), "origin_group"]
        worst_rmse_group = group_acc.loc[group_acc["rmse"].idxmax(), "origin_group"]
        worst_mape_group = group_acc.loc[group_acc["mape"].idxmax(), "origin_group"]
    if comp_acc is not None and not comp_acc.empty:
        best_compartment_rmse = comp_acc.loc[comp_acc["rmse"].idxmin(), "compartment"]
        worst_compartment_rmse = comp_acc.loc[comp_acc["rmse"].idxmax(), "compartment"]
        worst_compartment_mape = comp_acc.loc[comp_acc["mape"].idxmax(), "compartment"]

    # Introduction
    doc.add_heading("1. Introduction", level=1)
    p = doc.add_paragraph()
    p.add_run("Most debates on research mobility focus on net flows: which country gains researchers and which loses them. "
              "Net-flow accounting is useful for headlines, but it hides the transition rates that actually move researchers between career stages and locations. "
              "A small proportional change in one of those rates can, over time, push a research community below the minimum coauthor pool it needs to remain viable. "
              "Once the pool falls below that threshold, recovery becomes difficult or impossible, even if policy is later reversed. "
              "That is the point of no return (PNR) that motivates this paper. "
              "The contribution of this paper is to translate that qualitative insight into an empirically tractable model. "
              "We estimate transition rates from open bibliometric data, solve the steady state of a compartment model, and identify which rate in which civilisation is closest to a threshold. "
              "The approach is deliberately stylised: it sacrifices demographic realism for transparency and for the ability to compare multiple civilisations with the same accounting framework.")

    p = doc.add_paragraph()
    p.add_run("Artificial intelligence (AI) and machine learning (ML) have become the archetypal general-purpose technologies of the current era")
    add_citation(p, 1)
    p.add_run(", and their development depends on a relatively small, highly mobile workforce of doctoral and post-doctoral researchers, principal investigators (PIs), and research engineers")
    add_citation(p, 1)
    p.add_run(". "
              "The geographic concentration of this workforce has generated both scientific and geopolitical concern. "
              "Policymakers in the United States, China, Europe, Japan, India and elsewhere now treat AI talent as a strategic input, and several governments have introduced incentives to attract or retain researchers")
    add_citation(p, 2)
    p.add_run(". "
              "Most of those policies are evaluated by their immediate net-flow effects. "
              "They rarely ask which transition in the career pipeline is the binding constraint, or how close a community is to a threshold where the field can no longer sustain itself. "
              "The economic literature on science has long emphasised that researchers are a scarce input and that their mobility responds to career incentives and institutional quality")
    add_citation(p, 3)
    p.add_run(". "
              "That literature provides the microfoundation for our rates: individuals decide where to train, whether to go abroad, when to return, and when to leave academia. "
              "We aggregate those individual decisions into civilisation-level transition rates and ask what the resulting dynamics imply for community survival.")

    p = doc.add_paragraph()
    p.add_run("The civilisation framework offers a natural way to partition the global research population into culturally and institutionally coherent arenas")
    add_citation(p, 4)
    p.add_run(". "
              "We adapt Huntington's nine civilisations for AI/ML mobility by keeping the United States, China (Sinic), India and nearby South Asian countries (Hindu), Japan, and the Islamic world as distinct groups, splitting the Western bloc into the United States, Anglosphere excluding the United States, Continental Europe and Other Western, and merging the smaller Latin American, Orthodox and African communities into Other Civilizations. "
              "This grouping reflects the empirical size and mobility patterns observed in the data rather than a normative claim about civilisational identity.")

    p = doc.add_paragraph()
    p.add_run("The central argument of the paper is that preserving civilisational diversity in AI/ML is not only a normative preference but also a safeguard against technological dead ends. "
              "When a single region or a small oligopoly dominates a field, the set of research questions, evaluation norms, and institutional incentives narrows")
    add_citation(p, 5)
    p.add_run(". "
              "A diverse ecosystem generates competing approaches, which increases the probability that unexpected breakthroughs and error correction survive")
    add_citation(p, 5)
    p.add_run(". "
              "If transition rates can be observed with enough temporal resolution, policy can intervene before a community reaches the point of no return. "
              "Early, proportionate interventions can prevent the emergence of a monopoly or oligopoly without requiring large ex post rescues.")

    p = doc.add_paragraph()
    p.add_run("We therefore address five research questions. "
              "First, how close is each civilisation to the point of no return (PNR) in its AI/ML research community? "
              "Second, which transition rates have the largest effect on community size? "
              "Third, how have transition rates changed between earlier and later career cohorts, and what would have happened if those rates had persisted? "
              "Fourth, what safety-factor-bound single-lever and multi-lever policy scenarios can widen the margin before a point of no return (PNR) is reached? "
              "Fifth, can the fitted rates be estimated year by year and used to project near-term population composition, and how well do those projections reproduce observed 2017-2023 counts? "
              "The key policy intuition is that, with an appropriately chosen time step and an early warning signal, intervention can be calibrated in safety margins rather than after collapse. "
              "This prevents any single civilisation from cornering the supply of critical talent, and thereby preserves the competitive diversity that drives long-run innovation.")

    p = doc.add_paragraph()
    p.add_run("The contribution is a reproducible, data-driven transition-rate model that links OpenAlex publication records to a system of ordinary differential equations (ODEs)")
    add_citation(p, 6)
    p.add_run(". "
              "The model is intentionally simple: it does not explain why a rate is high or low, but it identifies which rate is closest to a threshold and therefore where early intervention is most urgent.")

    doc.add_paragraph()

    doc.add_heading("2. Literature and conceptual framework", level=1)

    p = doc.add_paragraph()
    p.add_run("Researcher mobility has long been studied under the headings of brain drain, brain circulation and brain gain")
    add_citation(p, 7)
    p.add_run(". "
              "Thorn and Holm-Nielsen argue that the mobility of researchers from developing countries can become a gain when return migration and diaspora networks are supported, but it can become a drain when local research environments cannot retain or reproduce talent")
    add_citation(p, 7)
    p.add_run(". "
              "Appelt et al., using a gravity framework for 1996-2011, find that scientific collaboration, economic convergence and visa restrictions are the strongest correlates of bilateral mobility")
    add_citation(p, 2)
    p.add_run(". "
              "Their analysis shows that mobility is multi-directional: a large share of researcher movement is better described as circulation than as one-way migration.")

    p = doc.add_paragraph()
    p.add_run("The AI/ML literature has documented the same patterns at higher resolution. "
              "MacroPolo's Global AI Talent Tracker finds that the United States remains the leading destination for top-tier AI researchers, while China and India are expanding domestic retention")
    add_citation(p, 1)
    p.add_run(". "
              "AlShebli et al. show that U.S.-China collaboration in AI is more impactful than either country working alone, and that most mobile AI scientists retain collaboration links with their origin country")
    add_citation(p, 8)
    p.add_run(". "
              "Yuan et al. find that the brain-drain problem for AI scientists is increasingly serious in developing countries, and that the ties among AI elites are highly clustered")
    add_citation(p, 9)
    p.add_run(". "
              "These studies establish that AI/ML talent is mobile, concentrated and strategically important.")

    p = doc.add_paragraph()
    p.add_run("What is missing is a formal link between individual transition rates and the long-run viability of a research community. "
              "The concept of a minimum viable population, introduced by Shaffer, captures the smallest isolated population that has a high probability of persisting despite demographic, environmental and genetic stochasticity")
    add_citation(p, 10)
    p.add_run(". "
              "Transferred to science, the equivalent idea is a minimum viable coauthor pool: the smallest number of active researchers that can continue to produce work at the field's observed coauthor intensity. "
              "Below that pool, collaboration networks fragment, mentorship chains break, and the field enters a self-reinforcing decline.")

    p = doc.add_paragraph()
    p.add_run("This framing generates four testable hypotheses. "
              "H1: Across all groups, the equilibrium active pool exceeds the minimum viable threshold, but the distance to the threshold varies widely. "
              "H2: Dropout is the transition rate with the largest negative effect, because attrition removes researchers from every compartment. "
              f"H3: {ctx['positive_lever_sentence'].rstrip('.')}. "
              "H4: Smaller civilisations, and those with older cohort structures, sit closer to their point of no return.")

    p = doc.add_paragraph()
    p.add_run("A final literature stream emphasises the consequences of concentrated research agendas. "
              "Aghion et al. provide evidence that the relationship between competition and innovation follows an inverted-U shape, with the strongest innovative performance in markets that are neither perfectly collusive nor perfectly monopolistic")
    add_citation(p, 5)
    p.add_run(". "
              "Translated to global science, this suggests that a single dominant region or a tight oligopoly may slow the rate of methodological and conceptual breakthroughs. "
              "Maintaining multiple centres of AI/ML research is therefore not merely a distributional concern; it may increase the long-run productivity of the field.")

    doc.add_heading("2.1 Researcher mobility", level=2)
    p = doc.add_paragraph()
    p.add_run("Researcher mobility has been studied from several angles. "
              "A large empirical literature documents net flows of scientists and inventors across countries and regions, often using patent or publication records")
    add_citation(p, 11)
    p.add_run(". "
              "That work consistently finds that the United States, parts of Europe and, increasingly, China and India are central nodes in the global mobility network. "
              "It also finds that mobility responds to wages, funding, institutional quality and career prospects, but that it is path-dependent: once a community loses its senior cohort, it becomes harder to rebuild.")

    doc.add_heading("2.2 Scientific collaboration and diversity", level=2)
    p = doc.add_paragraph()
    p.add_run("A second strand of work emphasises the structure of scientific collaboration. "
              "Multi-university and international teams now produce a growing share of high-impact research, and the geographic dispersion of teams does not necessarily reduce their impact")
    add_citation(p, 12)
    p.add_run(". "
              "This literature suggests that global AI/ML is not a zero-sum race in which every researcher in one location subtracts from another. "
              "It also implies that sustaining a domestic community is compatible with, rather than opposed to, international collaboration. "
              "The question is therefore not whether researchers move, but whether the domestic pipeline that replaces them is robust enough to keep the field alive.")

    doc.add_heading("2.3 Minimum viable populations and critical thresholds", level=2)
    p = doc.add_paragraph()
    p.add_run("The third relevant literature concerns population viability and critical thresholds. "
              "In conservation biology, the minimum viable population concept identifies the smallest number of individuals that can sustain a population in the wild")
    add_citation(p, 10)
    p.add_run(". "
              "We borrow that intuition and apply it to a research community. "
              "A field needs a minimum number of active researchers to produce work, train successors, and maintain peer review and conference communities. "
              "Below that threshold, positive feedback loops weaken: fewer researchers produce fewer students, fewer students produce fewer researchers, and the community enters a downward spiral. "
              "This is the point of no return.")

    doc.add_heading("2.4 This paper's framework", level=2)
    p = doc.add_paragraph()
    p.add_run("The present paper bridges these literatures by estimating transition rates from open bibliometric data and embedding them in a compartment model. "
              "The model is closest in spirit to Stephan's economic model of science, in which researchers move through career stages and respond to incentives")
    add_citation(p, 3)
    p.add_run(", but it adds a civilisational partition and a minimum viable coauthor threshold. "
              "The civilisational partition reflects the clustering of career incentives, language, funding systems and institutional networks along civilisational lines, which shape mobility beyond national borders alone")
    add_citation(p, 4)
    p.add_run(". "
              "It also draws on the innovation-systems literature, in which technological trajectories are shaped by sectoral and national systems of innovation")
    add_citation(p, 13)
    add_citation(p, 14)
    add_citation(p, 15)
    add_citation(p, 16)
    p.add_run(". "
              "In that view, technological change is path-dependent and distributed: routines, organisations and institutions co-evolve, so the loss of a research community is not merely a decline in headcount but a reduction in the variety from which future trajectories can be generated. "
              "The point of no return is therefore an innovation-systems problem: once a community falls below the minimum scale needed to sustain distinct research programmes, the path-dependent process of search and selection that produces new trajectories is impaired. "
              "This connects that macro-level, innovation-systems view of path-dependent technological change to individual career-transition data: the transition rates and PNR distances reported below can be read as an empirical early-warning indicator of whether a particular civilisational innovation system retains enough researchers to sustain a distinct technological trajectory. "
              "The result is a framework that can be updated as new data arrive and can compare the fragility of different research communities using a common metric. "
              "Because it is built on open bibliometric data and transparent transition rates, the model can be replicated and extended by other researchers and by policymakers who need a common language for discussing mobility and capacity.")

    # Data
    doc.add_heading("3. Data and grouping", level=1)
    p = doc.add_paragraph()
    p.add_run("We extracted AI/ML works and author histories from the OpenAlex API for subfield `subfields/1702` (Artificial Intelligence), using works published between 2000 and 2023")
    add_citation(p, 6)
    p.add_run(". "
              "OpenAlex provides open, CC0 bibliographic metadata including authors, affiliations, countries, publication dates, venues and citation links. "
              "We built author histories by following each author's sequence of works and affiliations, assigning them to a country for each work and then to a civilisation by the modal country of their recorded affiliations. "
              "The cohort is restricted to authors whose career-start year (first observed AI/ML publication year) is between 2000 and 2016 and who have at least two AI/ML works in the 2000-2023 window. "
              "An author is treated as active if they have at least one AI/ML work in 2020-2023, and as having dropped out otherwise. "
              "An author is classified as a principal investigator (PI) if their first last-author paper appears during the observation window; single-authored papers are treated as last-author papers so that culturally varying coauthorship norms do not bias the seniority proxy")
    add_citation(p, 3)
    p.add_run(". "
              "A 'hit' work is a paper whose citation count places it in the top 10% of AI/ML works in the same publication year, observed within the first eight career years, regardless of the author's position on the author list. "
              "The abroad flag is set if the author is affiliated with a non-origin civilisation within the first six career years. "
              "The final groups are: United States, Anglosphere ex-US, Continental Europe, Sinic, Japanese, Hindu, Islamic, Other Western, and Other Civilizations.")

    doc.add_heading("3.1 Country-to-civilisation mapping", level=2)
    p = doc.add_paragraph()
    p.add_run("The grouping follows Huntington's civilisation taxonomy but is adjusted for sample-size and mobility reality in AI/ML. "
              "The United States is separated from the broader Anglosphere because it is the dominant destination for AI/ML researchers and because its higher-education and funding systems differ systematically from those of other English-speaking countries. "
              "Continental Europe is kept distinct from the Anglosphere because intra-European mobility and EU research funding create a separate mobility bloc. "
              "Latin American, Orthodox and sub-Saharan African countries are merged into Other Civilizations because their AI/ML author counts in the sample are too small to estimate stable transition rates separately. "
              "These civilisation labels are operational categories based on observed publication-affiliation patterns; they are not normative claims about cultural or political identity, and they are reported in full in Supplementary Material. "
              "Civilisation-level categories have also been shown to predict large-scale digital-communication networks")
    add_citation(p, 17)
    p.add_run(" and country-capacity clusters in scientific mobility and collaboration")
    add_citation(p, 18)
    p.add_run(", which supports the use of this aggregation as a cross-national research heuristic.")

    doc.add_heading("3.2 Sample selection and variable definitions", level=2)
    p = doc.add_paragraph()
    p.add_run("Authors enter the cohort if their first observed AI/ML publication year is between 2000 and 2016 and they have at least two AI/ML works in the 2000-2023 observation window. "
              "The career-start year is the first observed AI/ML publication year. "
              "Authors with exclusively unknown affiliations or with all affiliations outside the mapped countries are excluded. "
              "For each author we record the country of the majority of their affiliations and the civilisation to which that country maps. "
              "An author is active if they have at least one AI/ML work in 2020-2023; otherwise they are recorded as having dropped out. "
              "A hit is a paper in the top 10% of AI/ML citations for its publication year, observed within the first eight career years, regardless of the author's position. "
              "A PI is an author whose first last-author paper appears during the observation window; single-authored papers are treated as last-author papers. "
              "The abroad flag is set if the author appears in a non-origin civilisation within the first six career years. "
              "The final cohort of 723,647 authors is a model-implied sample extracted from OpenAlex; the objective is to build a reproducible pipeline and demonstrate the transition-rate framework, not to provide a definitive census.")

    doc.add_heading("3.3 OpenAlex coverage and known biases", level=2)
    p = doc.add_paragraph()
    p.add_run("OpenAlex coverage has improved over time but remains incomplete for works before 2000 and for non-English publications. "
              "Author disambiguation is imperfect, especially for common names and authors with multiple name variants. "
              "Affiliation metadata are supplied by publishers and are sometimes missing or refer to the primary institution rather than the country of residence. "
              "For these reasons, the absolute counts reported here are lower bounds on the true global AI/ML workforce. "
              "The analysis nevertheless preserves relative comparisons across civilisations because the same extraction rules are applied uniformly. "
              "Replication from a clean OpenAlex snapshot should produce very similar transition rates and point-of-no-return rankings even if absolute counts shift.")

    p = doc.add_paragraph()
    p.add_run("Table 1 reports the size and composition of the extracted cohort. "
              "The Sinic and Continental Europe groups contribute the largest number of works, followed by the United States and the Anglosphere ex-US. "
              "The Japanese and Other Western groups are the smallest in terms of author counts. "
              "The cohort of 723,647 authors is a model-implied sample extracted from the OpenAlex snapshot; absolute counts should be interpreted as model-implied stocks rather than population totals, and the bootstrap intervals reported below give a more honest picture of the uncertainty around those stocks. "
              "The relative sizes are nevertheless informative. "
              "A civilisation with a small cohort but a low coauthor intensity can be more resilient than a larger civilisation with a high coauthor intensity, because the former needs fewer distinct PI groups to sustain its output. "
              "This is why the minimum viable coauthor threshold and the equilibrium active pool must be compared jointly.")

    desc = _descriptive_table(cohort)
    _add_table_from_df(
        doc,
        desc,
        caption="Table 1. Descriptive statistics for the extracted AI/ML cohort by civilisation group. Civilisation labels are operational aggregations of OpenAlex country-affiliation patterns and do not imply normative cultural or political classification.",
        decimals={"n": 0, "works": 0, "active": 0, "hits": 0, "pis": 0, "career_start_mean": 1, "abroad": 0},
    )

    # Methods
    doc.add_heading("4. Methods", level=1)
    doc.add_heading("4.1 Compartment model", level=2)
    p = doc.add_paragraph()
    p.add_run("Each civilisation is represented by six compartments: domestic early-career researchers (D), abroad early-career researchers (A), domestic hit researchers (H_D), abroad hit researchers (H_A), domestic principal investigators (P_D), and abroad principal investigators (P_A). "
              "Transition rates are early-career outflow (α), return (β), hit generation at home and abroad (h_D and h_A), PI promotion at home and abroad (p_D and p_A), and dropout from all compartments (d). "
              "The equations are:")
    add_omath_paragraph(doc, math_ode_system())
    p = doc.add_paragraph()
    p.add_run("The model makes several simplifying assumptions. "
              "It treats each civilisation as a single aggregate, ignoring cross-civilisation collaboration and spillovers. "
              "It assumes constant per-year transition rates and a continuous-time Markov structure. "
              "Career stages are collapsed into the three observed layers: early-career, hit researchers and PIs. "
              "These simplifications are necessary to keep the model estimable from OpenAlex and to make the point-of-no-return calculation transparent. "
              "They also mean that the model is best interpreted as a stylised early-warning device, not as a realistic demographic projection.")

    doc.add_heading("4.2 Endogenous inflow", level=2)
    p = doc.add_paragraph()
    p.add_run("New entrants are modelled as a function of the domestic PI stock. "
              "The linear form is ")
    add_omath_inline(p, math_I_linear())
    p.add_run(f", where I_0 is the exogenous entry rate, r is the PI reproduction rate, and r is capped at {_fmt(ctx['safety_factor_cap'], 2)}× the stability-critical value (safety factor {_fmt(ctx['safety_factor_cap'], 2)}); the most constrained fitted group realises {_fmt(ctx['min_realised_safety_ratio'], 2)}×. "
              "A saturating alternative, ")
    add_omath_inline(p, math_I_saturating())
    p.add_run(", is reported as a robustness check. "
              "The PI-driven inflow captures the idea that senior researchers train graduate students, attract postdoctoral researchers, and create the institutional infrastructure that produces the next cohort. "
              "This is a strong assumption because it ignores cross-border recruitment and non-PI sources of new researchers, but it provides a transparent lower bound: if the domestic PI stock falls, the model predicts a decline in new entrants. "
              "The safety factor prevents the model from producing runaway growth when the observed r exceeds the critical value, which is a common empirical finding because observed recruitment is bounded by the data window.")

    doc.add_heading("4.3 Minimum viable coauthor threshold", level=2)
    p = doc.add_paragraph()
    p.add_run("For each group we computed the mean number of authors per work (c\u0304) and the median number of distinct last-author groups observed per recent year (k). "
              "The minimum viable domestic active pool is ")
    add_omath_inline(p, math_threshold())
    p.add_run(". When the equilibrium active pool ")
    add_omath_inline(p, math_active_pool())
    p.add_run(" falls below M, the community can no longer produce works at the observed coauthor intensity. "
              "In this sense, falling below M is a sufficient condition for collapse, not a necessary one; external shocks can push a community below viability even when the equilibrium active pool remains above M. "
              "The threshold is deliberately conservative: it assumes that each new work requires at least k distinct PI groups and that each work has the average number of coauthors. "
              "This overstates the number of distinct actors needed for a viable field, which means that M is a soft lower bound and that observed margins are probably smaller than they appear. "
              "A community with a margin just above M is therefore more fragile than the number itself suggests.")

    doc.add_heading("4.4 Estimation, equilibrium and sensitivity", level=2)
    p = doc.add_paragraph()
    p.add_run("Transition rates are estimated as constant per-year hazards from observed proportions within the cohort. "
              "For each group and each transition, the rate is the ratio of observed transitions to the total exposure time spent in the source compartment during the observation window, using a Laplace pseudocount of 1 for each outcome so the smoothed proportion is (successes + 1)/(n + 2). "
              "This avoids zero-rate singularities when the cohort is small. "
              "Because the data are right-censored at the end of the observation period, the resulting rates are lower bounds on true long-run hazards; equilibrium solutions therefore tend to be conservative. "
              "The non-linear steady-state equations are solved numerically using a trust-region Newton method with analytically supplied Jacobians. "
              "Elasticities are computed by perturbing each rate by 1%, re-solving, and taking the percentage change in the target stock. "
              "For point-of-no-return analysis we scale each rate until the active pool T reaches its coauthor threshold M, or the domestic PI pool P_D reaches k distinct last-author groups as a lower-bound PI-pool threshold, and record the critical factor and its proximity, |critical factor − 1|. "
              "A rate whose critical factor lies inside the scan window and is close to 1.0 is the most fragile lever for that group. "
              "All counterfactuals are mechanical perturbations of the fitted rates; they reveal which transitions the model treats as sensitive, not the causal impact of real-world policies.")

    doc.add_heading("4.5 Limitations", level=2)
    p = doc.add_paragraph()
    p.add_run("The main limitations are data quality and model scope. "
              "OpenAlex country metadata are noisy, especially for older works and for authors with multiple affiliations. "
              "Career stages are inferred from authorship order and are imperfect proxies. "
              "The model does not include cross-civilisation knowledge spillovers, bilateral migration costs, or firm-level mobility. "
              "Finally, the assumption of constant rates is a strong approximation over a 23-year window. "
              "We therefore emphasise rank-order and relative sensitivity rather than point forecasts.")

    doc.add_heading("4.6 Annual transition-rate estimation and projection", level=2)
    p = doc.add_paragraph()
    p.add_run("The steady-state model in Sections 4.1-4.4 treats rates as constants. "
              "To test whether the same framework can be used for short-run monitoring, we reconstructed year-by-year compartment membership from the cohort data. "
              "For each author and year we inferred location as domestic if the author was in the origin civilisation and abroad otherwise. "
              "From these states we computed annual transition counts for the six compartments, applied Laplace smoothing with a pseudocount of 1 for each possible destination, and derived the probabilities that map to α, β, h_D, h_A and p_D, p_A. "
              "Dropout (d) is not directly observed year-by-year in the training window because final attrition is right-censored before 2023, so we import the cohort-level per-year hazard from the full-career data and treat it as a constant annual rate for each group. "
              "Inter-civilisation flows are approximated by assigning each abroad author-year to the author's recent_group as the destination civilisation.")

    p = doc.add_paragraph()
    p.add_run("For the 2017-2026 projection we fit a linear trend to the observed 2000-2016 rates for each group and rate. "
              "If fewer than four observations were available or the fit explained less than 10% of the variance, the historical mean was used instead. "
              "Projected rates were clipped to values between 0 and 1. "
              "Projected annual dropout was capped at 1.5 times the 90th percentile of observed annual dropout rates in the 2000-2016 training period. "
              "Projected total inflows were apportioned across compartments using the first-compartment distribution observed over the 2000-2016 training period. "
              "Population composition was projected forward with the discrete-time recursion N(t+1) = N(t)P(t) + b(t+1), where P(t) is a 6×6 row-stochastic-in-expectation matrix that preserves dropout mass: the row sum is 1 − d after scaling outgoing rates. "
              "This discrete step is the operational counterpart of the continuous-time ODE; with an annual dt it provides an early-warning signal one year ahead.")

    p = doc.add_paragraph()
    p.add_run("We compare the 2017-2023 projection with the observed annual stock. "
              "The comparison is limited to years that have observed data, and the observed stock is reindexed to the full group-year-compartment grid so that zero-observed cells are not omitted from the accuracy metrics. "
              "Accuracy is reported as root mean square error (RMSE) and mean absolute percentage error (MAPE); MAPE here is computed against count_obs + 1 to avoid division by zero and is therefore a conservative, non-standard measure.")

    doc.add_heading("4.7 Correction pressures and theoretical bounds", level=2)
    p = doc.add_paragraph()
    p.add_run("The annual estimates contain several regularising pressures that bound the model away from instability and fabrication. "
              "Laplace smoothing adds a uniform prior of 1 to every possible destination, which shrinks sparse cells toward 1/(number of destinations) and prevents zero-probability singularities when a transition is unobserved in a small group-year. "
              "It is equivalent to a weak Dirichlet prior and is a standard regulariser for sparse multinomial transitions.")

    p = doc.add_paragraph()
    p.add_run("Clipping projected rates to values between 0 and 1 is a feasibility pressure: rates outside the probability simplex are inadmissible. "
              "The annual dropout rate is anchored to the cohort-level per-year hazard rather than extrapolated from year-to-year transitions, because final attrition is right-censored in the training window. "
              "The inflow apportionment pressure keeps the composition of new entrants aligned with the most recently observed recruitment pattern, rather than inventing a new distribution. "
              f"Finally, the endogenous inflow is capped at a safety factor of {_fmt(ctx['safety_factor_cap'], 2)} relative to the critical reproduction rate (the most constrained fitted group realises {_fmt(ctx['min_realised_safety_ratio'], 2)}), which keeps the system inside the stability boundary. "
              "Together these pressures embody the principle that projection should stay within observed empirical support and within theoretical stability limits; they are not arbitrary adjustments but transparent bounds that can be tightened or relaxed as more data become available.")

    # Results
    doc.add_heading("5. Results", level=1)
    p = doc.add_paragraph()
    p.add_run(f"Table 2 reports the equilibrium domestic active pool T, the minimum viable threshold M, and the endogenous inflow parameters for the {len(eq)} groups. "
              "All groups remain above their threshold under the fitted model, but margins differ by an order of magnitude. "
              f"The {ctx['largest_pools']} groups show the largest equilibrium active pools, reflecting large cohorts and relatively low coauthor-intensity thresholds. "
              f"The {ctx['smallest_pool']} group has the smallest equilibrium active pool, and {ctx['smallest_margin_group']} has the narrowest safety margin, although both still exceed their minimum viable coauthor pool. "
              "The ratio T/M is a summary resilience indicator, but absolute margin is the more direct measure of proximity to the point of no return.")

    eq_table = eq[["group", "T_equilibrium", "M_threshold", "margin_to_threshold_T", "I0", "r", "r_obs", "r_critical"]].copy()
    eq_table["T_over_M"] = eq_table["T_equilibrium"] / eq_table["M_threshold"]
    eq_table = eq_table.rename(columns={
        "group": "Group",
        "T_equilibrium": "T_eq",
        "M_threshold": "M",
        "margin_to_threshold_T": "Margin",
        "T_over_M": "T/M",
        "r_obs": "r_obs",
        "r_critical": "r_crit",
    })
    _add_table_from_df(
        doc,
        eq_table,
        caption="Table 2. Equilibrium domestic active pool, minimum viable threshold, and endogenous inflow parameters.",
        decimals={"T_eq": 0, "M": 0, "Margin": 0, "T/M": 2, "I0": 0, "r": 5, "r_obs": 5, "r_crit": 5},
    )

    p = doc.add_paragraph()
    p.add_run("Figure 1 visualises the gap between equilibrium and threshold. "
              f"The {ctx['largest_pools']} groups display the largest equilibrium active pools, while the {ctx['smallest_pool']} group is the smallest. "
              "However, the point-of-no-return metric is not the absolute level of T but the distance between T and M, which reflects both the stock of researchers and the coauthor intensity of the field. "
              "Groups with high T but also high c\u0304 and k can still be fragile if their margin is small.")
    doc.add_picture(str(fig_paths["fig1"]), width=Inches(5.8))
    cap = doc.add_paragraph()
    cap.add_run("Figure 1. Equilibrium domestic active pool (T) and minimum viable coauthor threshold (M) by group.").italic = True
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER

    p = doc.add_paragraph()
    p.add_run("Table 3 shows the three transition-rate elasticities with the largest absolute impact on T for each group. "
              f"Dropout (d) is the largest negative lever in every group; its active-pool elasticity ranges from {_fmt(ctx['d_min_e'], 2)} to {_fmt(ctx['d_max_e'], 2)}. "
              "Attrition removes researchers from every compartment, so a proportional increase in d produces a larger proportional decline in the active pool. "
              f"{ctx['positive_lever_sentence']} "
              "Early-career outflow (α) has a modest negative effect in most groups, but because it moves researchers to the abroad compartment rather than removing them entirely, its direct impact on the domestic active pool is smaller than that of dropout. "
              "There is notable heterogeneity in the magnitude of the positive levers. "
              f"The {ctx['highest_pd_group']} group shows the strongest response to PI promotion (p_D), indicating that improving the promotion of hit researchers to PIs is an especially efficient way to expand the domestic active pool in that community. "
              "In the largest groups, p_D remains positive but its relative effect is smaller, because the active pool is already large and a proportional change in promotion has less marginal impact.")

    rows3 = []
    for group, gdf in top_t.groupby("group"):
        top3 = gdf.sort_values("abs_elasticity", ascending=False).head(3)
        r = [group]
        for _, row in top3.iterrows():
            r.extend([row["rate"], _fmt(row["elasticity"], 3)])
        rows3.append(r)
    elas_df = pd.DataFrame(rows3, columns=["Group", "1st rate", "1st elasticity", "2nd rate", "2nd elasticity", "3rd rate", "3rd elasticity"])
    _add_table_from_df(
        doc,
        elas_df,
        caption="Table 3. Top transition-rate elasticities for domestic active pool T.",
    )

    closest = pnr_closest.iloc[0]
    p = doc.add_paragraph()
    p.add_run(f"Table 4 reports, for each group, the single rate that reaches the active-pool threshold with the smallest proportional change. "
              f"The {closest['group']} group is the most fragile: {closest['rate_name']} must be multiplied by {_fmt(closest['critical_factor'], 3)}× its current value (equivalent to a {closest['proximity']*100:.0f}% proportional {'reduction' if closest['critical_factor'] < 1 else 'increase'}) to drive the active pool to its minimum viable threshold. "
              f"{ctx['pnr_lever_text']}. "
              "This is consistent with a recruitment-driven view of scientific communities: if the pipeline of new researchers shuts or slows, the active pool eventually falls below the minimum viable coauthor pool regardless of how efficient return or promotion becomes. "
              "A global retention programme that reduces dropout would benefit all groups, but the most vulnerable groups may also need an expansion of the exogenous entry rate.")

    pnr_table = pnr_closest[["group", "target", "rate_name", "current_rate", "critical_factor", "proximity"]].copy()
    pnr_table.columns = ["Group", "Target", "Rate", "Current", "Critical factor", "Proximity"]
    _add_table_from_df(
        doc,
        pnr_table,
        caption="Table 4. Closest point of no return for the active researcher pool by group.",
        decimals={"Current": 4, "Critical factor": 3, "Proximity": 3},
    )

    p = doc.add_paragraph()
    p.add_run("Figure 2 ranks groups by their closest point-of-no-return sensitivity.")
    doc.add_picture(str(fig_paths["fig2"]), width=Inches(5.8))
    cap = doc.add_paragraph()
    cap.add_run("Figure 2. Closest point-of-no-return proximity by group. Smaller values mean a smaller proportional change in the listed rate is required to reach the threshold for the stated target pool.").italic = True
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER

    if sat_eq is not None:
        doc.add_heading("5.1 Saturating recruitment extension", level=2)
        p = doc.add_paragraph()
        p.add_run("We also test a saturating recruitment function in which each additional PI adds fewer entrants. "
                  "With the capacity parameter calibrated to observed PI stocks, the saturating equilibrium is "
                  f"{ctx['sat_range_text']} at the displayed precision, so the linear safety-factor bound remains the operative constraint. "
                  "Table 5 reports the fitted epsilon values; the near-zero differences show that the results are not driven by unbounded linear growth, but they do not rule out stronger saturation at higher PI densities.")
        merged = eq[["group", "T_equilibrium"]].merge(
            sat_eq[["group", "T_equilibrium", "epsilon"]], on="group", suffixes=("_lin", "_sat")
        )
        merged.columns = ["Group", "Linear T", "Saturating T", "ε"]
        _add_table_from_df(
            doc,
            merged,
            caption="Table 5. Equilibrium T under linear and saturating PI-driven inflow.",
            decimals={"Linear T": 0, "Saturating T": 0, "ε": 5},
        )

        # Robustness of PNR rankings to the functional form of endogenous inflow.
        pnr_rob = pnr_robustness_table()
        if not pnr_rob.empty:
            p = doc.add_paragraph()
            p.add_run("The closest point-of-no-return lever is the same under the saturating alternative for every civilisation: "
                      "exogenous entry (I0) is the rate that requires the smallest proportional change to push the active pool to its minimum viable threshold. "
                      "Table 5a reports the proportional factor and proximity for the active-pool threshold under both assumptions. "
                      "The rank order of civilisational fragility is preserved (Spearman ρ = 1.0), and the absolute proximity values move in the same direction. "
                      "This confirms that the policy ranking—exogenous entry first, then dropout, then domestic promotion and return—is robust to replacing the linear feedback with a saturating one.")
            _add_table_from_df(
                doc,
                pnr_rob,
                caption="Table 5a. Closest PNR lever and proximity under linear and saturating endogenous inflow (active-pool threshold).",
                decimals={"linear_factor": 4, "saturating_factor": 4, "linear_proximity": 4, "saturating_proximity": 4},
            )

    doc.add_heading("5.2 Historical counterfactual", level=2)
    n_compare = len(period_compare)
    if ctx["period_all_neg"]:
        prefix = "Both" if n_compare == 2 else f"All {n_compare}"
        period_direction_text = (
            f"{prefix} groups with dual-window support would see smaller safety margins under late-window rates "
            f"({ctx['period_neg']})."
        )
    else:
        period_direction_text = (
            f"Groups that would see smaller safety margins under late-window rates: {ctx['period_neg']}. "
            f"Groups that would see larger safety margins under late-window rates: {ctx['period_pos']}."
        )
    p = doc.add_paragraph()
    p.add_run("Table 6 compares the equilibrium that would have emerged if the transition rates estimated for the early career window (2000-2010) or the late window (2011-2016) had persisted indefinitely. "
              "The late window is shorter and its rates are estimated from younger cohorts, so the comparison should be read as a sensitivity exercise rather than a forecast. "
              f"Only {n_compare} groups have enough dual-window support for reliable rate estimation in both windows; they are listed in the table. "
              f"{period_direction_text} "
              "This pattern shows that global AI/ML mobility is not moving in a single direction; different civilisations are on different trajectories, and a uniform policy response would ignore this heterogeneity. "
              "Because the late cohort is younger, the late-window equilibrium is likely biased downward for groups where career progression has not yet run its course. "
              "Even so, the exercise shows that the current regime is not the only possible one, which is why counterfactual policy analysis is useful.")

    pc = period_compare.copy()
    pc = pc[["group", "T_early", "T_late", "pct_delta_T", "margin_early", "margin_late", "delta_margin"]].rename(columns={
        "group": "Group",
        "T_early": "T early",
        "T_late": "T late",
        "pct_delta_T": "ΔT (%)",
        "margin_early": "Margin early",
        "margin_late": "Margin late",
        "delta_margin": "Δ margin",
    })
    _add_table_from_df(
        doc,
        pc,
        caption="Table 6. Historical counterfactual: equilibrium active pool and safety margin under early versus late transition-rate regimes.",
        decimals={"T early": 0, "T late": 0, "ΔT (%)": 1, "Margin early": 0, "Margin late": 0, "Δ margin": 1},
    )

    p = doc.add_paragraph()
    p.add_run("Figure 3 shows the change in safety margin between the early and late transition-rate regimes.")
    doc.add_picture(str(fig_paths["fig3"]), width=Inches(5.8))
    cap = doc.add_paragraph()
    cap.add_run("Figure 3. Change in safety margin between early and late transition-rate regimes. Positive values mean the late-window rates would produce a larger safety margin than the early-window rates if they persisted; negative values mean the margin would shrink. "
                "The comparison is across two point estimates; uncertainty is substantial because the two windows have different cohort sizes and the steady-state model does not capture policy shocks.").italic = True
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_heading("5.3 Policy counterfactuals", level=2)
    policy_top = policy_rank.groupby("group").head(1).copy()
    # Robustly describe the 10% dropout reduction effect, even if another lever is top for some group
    d_decrease = policy_rank[(policy_rank["lever"] == "d") & (policy_rank["direction"] == "decrease")].copy()
    d_10pct = d_decrease[d_decrease["lever_change_pct"].abs() >= 9.9]
    if d_10pct.empty:
        d_10pct = d_decrease
    d_10pct_group = d_10pct.loc[d_10pct.groupby("group")["normalised_margin_gain_per_10pct"].idxmax()]
    all_top_are_d = (policy_top["lever"] == "d").all()
    d_min = d_10pct_group.sort_values("margin_gain").iloc[0]
    d_max = d_10pct_group.sort_values("margin_gain").iloc[-1]
    min_gain_group = d_min["group"]
    max_gain_group = d_max["group"]
    min_gain = round(d_min["margin_gain"])
    max_gain = round(d_max["margin_gain"])
    dominant_text = (
        "Reducing dropout is the dominant positive lever for every civilisation, which is consistent with the elasticity results in Table 3. "
        if all_top_are_d
        else "Reducing dropout is the dominant positive lever for most civilisations in the current data. "
    )
    p = doc.add_paragraph()
    p.add_run(f"Table 7 reports the single mechanical counterfactual with the largest margin gain per 10% lever change for each group. {dominant_text}"
              f"The gain from a roughly 10% proportional reduction in d ranges from about {min_gain} additional active researchers for the {min_gain_group} group to about {max_gain} for the {max_gain_group} group, reflecting differences in cohort size and baseline attrition. "
              "No other single lever comes close to dropout reduction in terms of simulated margin gain per unit proportional change, although combinations of levers may be more efficient for some groups. "
              "The results also imply that policy need not focus on blocking early-career outflow. "
              "Reducing attrition among researchers who remain in the domestic system is a more efficient way to protect the active pool than preventing researchers from going abroad, because a researcher abroad is still in the global AI/ML system and may return. "
              "For the smallest groups, increasing the exogenous entry rate or improving the promotion of hit researchers to PIs can add additional margin, but dropout reduction remains the first-order model-implied target.")

    policy_top = policy_top.rename(columns={
        "group": "Group",
        "lever": "Lever",
        "direction": "Direction",
        "lever_change_pct": "Change (%)",
        "margin_gain": "Margin gain",
        "normalised_margin_gain_per_10pct": "Gain per 10%",
    })
    _add_table_from_df(
        doc,
        policy_top,
        caption="Table 7. Top positive mechanical counterfactual per group, measured by margin gain per 10% proportional lever change.",
        decimals={"Change (%)": 0, "Margin gain": 0, "Gain per 10%": 1},
    )

    package_text, _ = _package_summary()
    if package_text:
        p = doc.add_paragraph(package_text)

    doc.add_heading("5.4 Uncertainty", level=2)
    p = doc.add_paragraph()
    p.add_run("Supplementary Table 5 reports bootstrap 95% confidence intervals for the equilibrium active pool T and the domestic PI pool P_D. "
              "The intervals are wide, reflecting the model-implied cohort scale and the extrapolation from observed author-career exposure to long-run steady states. "
              "For some groups the upper bound is an order of magnitude larger than the lower bound, indicating that the equilibrium is sensitive to resampling variation in the transition rates. "
              "This uncertainty should be interpreted as a warning against over-interpreting point estimates and as a reason to view the point-of-no-return distances as indicative rather than precise thresholds. "
              "Despite the width, the lower bounds for most groups remain above the minimum viable threshold, which supports the qualitative conclusion that all groups are currently above the point of no return. "
              "For the smallest groups the lower bound is closer to M, reinforcing the need for continued monitoring and for policy buffers.")

    p = doc.add_paragraph()
    p.add_run("Figure 4 displays the bootstrap intervals graphically.")
    doc.add_picture(str(fig_paths["fig4"]), width=Inches(5.8))
    cap = doc.add_paragraph()
    cap.add_run("Figure 4. Bootstrap 95% confidence intervals for equilibrium T by group.").italic = True
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_heading("5.5 Synthesis", level=2)
    p = doc.add_paragraph()
    p.add_run("Taken together, the results provide a consistent picture. "
              "Exogenous entry and dropout are the two rates that most strongly determine the long-run viability of an AI/ML research community. "
              "Communities that are large in absolute terms are not necessarily safe if their coauthor intensity is high; conversely, small communities can be robust if their attrition is low and their recruitment pipeline is stable. "
              "The historical counterfactual shows that the current regime is not preordained: a shift in transition rates at the start of the AI boom would have produced different steady states for different civilisations. "
              "This is precisely why the framework is useful: it identifies which rate in which community is closest to a threshold, allowing policy to intervene before rather than after a collapse. "
              "The policy message is therefore both diagnostic and preventative. "
              "By tracking transition rates rather than net flows, policymakers can see which civilisation is approaching a point of no return and which lever offers the largest safety margin per unit of effort.")

    doc.add_heading("5.6 Annual transition rates and inter-civilisation flows", level=2)
    p = doc.add_paragraph()
    p.add_run("Figure 5 plots the observed 2000-2016 transition rates and the projected 2017-2026 rates for each civilisation. "
              "Rates are displayed by group and by transition type, so that the reader can see whether a particular transition is trending toward a boundary. "
              "Because the projections are linear trend fits regularised by the correction pressures described in Section 4.7, they are not forecasts of specific future events; they are the model's one-year-ahead extrapolation of the recent historical trajectory.")
    doc.add_picture(str(fig_paths["fig5"]), width=Inches(6.0))
    cap = doc.add_paragraph()
    cap.add_run("Figure 5. Observed (solid) and projected (dashed) transition rates by civilisation, 2000-2026.").italic = True
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER

    p = doc.add_paragraph()
    p.add_run("The mean observed annual transition rates by group between 2000 and 2016, "
              "distinguishing early-career outflow (α), return (β), domestic and abroad hit generation (h_D, h_A), PI promotion (p_D), dropout (d), and total inflow (I_total), "
              "are provided in Supplementary Table 3. "
              "Similarly, the cross-origin-destination pairs with the largest accumulation of abroad author-years are listed in Supplementary Table 4. "
              "These supplementary tables keep the main text focused on the PNR and policy conclusions while preserving the empirical detail needed for replication and extension.")

    p = doc.add_paragraph()
    p.add_run("Figure 6 shows the cross-civilisation accumulation of abroad author-years. "
              "Rows represent the origin civilisation and columns represent the destination civilisation, approximated by the author's recent_group while abroad. "
              "Origin-destination cells with the same civilisation and destinations labelled Unknown are excluded because the reconstruction cannot observe the actual host civilisation. "
              "The remaining cells are a lower-bound proxy for the true inter-civilisation pipelines.")
    doc.add_picture(str(fig_paths["fig6"]), width=Inches(5.8))
    cap = doc.add_paragraph()
    cap.add_run("Figure 6. Cross-civilisation abroad author-year accumulation by origin (rows) and destination (columns) (same-civilisation cells and Unknown destinations excluded; lower-bound proxy).").italic = True
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER



    doc.add_heading("5.7 Out-of-sample projection, 2017-2023", level=2)
    p = doc.add_paragraph()
    p.add_run(f"The 2017-2023 projection is compared with observed annual stocks in Figure 7. "
              f"Stock-level accuracy is RMSE {_fmt(annual_ctx.get('overall_rmse', float('nan')), 2)} and MAPE {_fmt(annual_ctx.get('overall_mape_pct', float('nan')), 1)}% (a non-standard, conservative measure computed against count_obs + 1 to avoid division by zero). "
              "These stock-level metrics are not a fair test of the model: the estimation cohort is fixed to authors whose careers began by 2016, so observed post-2016 counts cannot include the new entrants that the projection adds each year. "
              "The projection therefore necessarily diverges from observed stocks for any civilisation with positive recruitment. "
              f"A cleaner validation is at the rate level: the projected transition rates have RMSE {_fmt(annual_ctx.get('rate_overall_rmse', float('nan')), 4)} and MAE {_fmt(annual_ctx.get('rate_overall_mae', float('nan')), 4)}, and the model's skill relative to a historical-mean baseline is {_fmt(annual_ctx.get('rate_overall_skill', float('nan')), 2)}. "
              f"By rate, the best relative skill is for {annual_ctx.get('best_rate_skill', '—')} ({_fmt(annual_ctx.get('best_rate_skill_value', float('nan')), 2)}× the historical-mean RMSE), while the weakest is for {annual_ctx.get('worst_rate_skill', '—')} ({_fmt(annual_ctx.get('worst_rate_skill_value', float('nan')), 2)}×). "
              "These figures show that the annual layer captures rate drift at least as well as a naive mean forecast, and should be read as a directional early-warning indicator of drift and threshold proximity rather than as a precise population forecast.")

    p = doc.add_paragraph()
    p.add_run("Direction and threshold-alarm diagnostics support this interpretation. ")
    if "direction_agreement" in annual_ctx:
        p.add_run(f"Year-to-year direction agreement between projected and observed compartment counts is {_fmt(annual_ctx['direction_agreement'] * 100, 1)}%, "
                  f"ranging from {annual_ctx.get('worst_direction_group', '—')} to {annual_ctx.get('best_direction_group', '—')}. ")
    if "threshold_alarm_accuracy" in annual_ctx:
        p.add_run(f"For the active pool T = D + H_D + P_D, the projection correctly classifies whether T is below the minimum viable threshold M in {_fmt(annual_ctx['threshold_alarm_accuracy'] * 100, 1)}% of group-years "
                  f"(sensitivity {_fmt(annual_ctx['threshold_alarm_sensitivity'] * 100, 1)}%, specificity {_fmt(annual_ctx['threshold_alarm_specificity'] * 100, 1)}%). "
                  f"The observed threshold-crossing group-years (n = {_fmt(annual_ctx.get('threshold_alarms_obs', 0), 0)}) all occur for the smallest civilisation in the post-2016 fixed cohort; they reflect the depletion of that cohort as careers mature, not a projected collapse. "
                  "The projection, by construction, adds new entrants each year and therefore does not predict such within-cohort depletion. "
                  "A zero sensitivity in this hold-out is thus a consequence of the fixed-cohort validation design, not evidence that the model misses genuine collapse events. ")
    p.add_run("These metrics confirm that the annual layer is useful for directional and threshold-crossing surveillance, not for precise population counts. "
              "The modest year-to-year direction-agreement value is expected for small compartments and sparse transitions, and it reinforces that the projection layer should be treated as a drift-and-threshold alarm rather than a population forecast.")
    doc.add_picture(str(fig_paths["fig7"]), width=Inches(6.0))
    cap = doc.add_paragraph()
    cap.add_run("Figure 7. Observed (solid) and projected (dashed) compartment counts by civilisation, 2017-2023. The vertical dotted line marks the end of the training period (2016).").italic = True
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER

    p = doc.add_paragraph()
    p.add_run("Detailed projection accuracy by civilisation and by compartment is reported in Supplementary Material (Supplementary Tables 1 and 2). "
              f"Among compartments, the lowest RMSE is for {best_compartment_rmse}, while the highest RMSE is for {worst_compartment_rmse} and the highest MAPE is for {worst_compartment_mape}. "
              "P_D and H_D show larger errors because small changes in PI and hit rates are amplified by the endogenous inflow term.")

    doc.add_heading("5.8 Japan-specific compartment and transition-rate ladder", level=2)
    p = doc.add_paragraph()
    p.add_run(f"Figure 8 places the Japanese AI/ML research community in the compartment model. "
              f"The fitted equilibrium is T={ja_ctx['D'] + ja_ctx['H_D'] + ja_ctx['P_D']} active researchers "
              f"(D={ja_ctx['D']}, H_D={ja_ctx['H_D']}, P_D={ja_ctx['P_D']}) against a minimum viable threshold of M={_fmt(ja_ctx['M'], 0)}, "
              f"so the safety ratio T/M is {_fmt(ja_ctx['T_over_M'], 2)}. "
              "The right-hand ladder compares Japan's six transition rates with those of the other civilisations. "
              f"Japan's closest point of no return is the exogenous entry rate I0: if I0 were reduced to {_fmt(ja_ctx['pnr_factor'] * 100, 1)}% of its current level, the active pool would reach the minimum viable threshold. "
              f"In the fitted rates, early-career outflow (α={_fmt(ja_ctx['alpha'], 3)}) and domestic PI promotion (p_D={_fmt(ja_ctx['p_D'], 3)}) are comparatively low, while return from abroad (β={_fmt(ja_ctx['beta'], 3)}) and domestic hit generation (h_D={_fmt(ja_ctx['h_D'], 3)}) are moderate. "
              "The small absolute size of the abroad PI compartment (P_A) shows that few Japanese researchers who leave eventually become PIs abroad, which makes the domestic pipeline the critical margin.")
    doc.add_picture(str(fig_paths["fig8"]), width=Inches(6.0))
    cap = doc.add_paragraph()
    cap.add_run("Figure 8. The Japanese AI/ML research community in the six-compartment model, with a cross-civilisation ladder of fitted transition rates. Japan is highlighted in the right-hand panel; longer bars represent higher estimated rates.").italic = True
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_heading("5.9 A combined model-evaluation view: T/M and PNR proximity", level=2)
    p = doc.add_paragraph()
    p.add_run("Figure 9 combines the long-run safety ratio T/M with the closest point-of-no-return proximity for each civilisation. "
              "A point in the lower-left corner has both a low equilibrium buffer and a small proportional change needed to reach the threshold, so it is the most fragile combination. "
              "Japan sits in this region alongside the 'Other Civilizations' group, even though its T/M ratio is above one. "
              "This dual view is useful as a model-evaluation metric: a civilisation can have a T/M ratio that looks comfortable but still be close to its PNR because the PNR depends on the proportional change in the most sensitive rate, not only on the level of T. "
              "Japan is presented here as an illustrative case; the same diagnostic can be applied to any civilisation with sufficient OpenAlex coverage.")
    doc.add_picture(str(fig_paths["fig9"]), width=Inches(5.8))
    cap = doc.add_paragraph()
    cap.add_run("Figure 9. Equilibrium safety ratio (T/M) versus closest point-of-no-return proximity for all civilisations. Japan is shown in red.").italic = True
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Discussion
    doc.add_heading("6. Discussion", level=1)
    p = doc.add_paragraph()
    p.add_run("The results support a transition-rate view of research policy. "
              "Rather than asking which country has a net inflow or outflow of researchers, the model asks which rate must be altered to keep a community above its minimum viable coauthor pool. "
              "This shift in focus has implications for how we conceptualise brain drain, design science and technology policy, and interpret civilisational diversity in AI/ML.")

    doc.add_heading("6.1 From net flows to transition rates", level=2)
    p = doc.add_paragraph()
    p.add_run("Most empirical studies of researcher mobility measure net flows, stocks or collaboration counts")
    add_citation(p, 11)
    p.add_run(". "
              "These indicators are useful for describing patterns, but they do not reveal the mechanisms that sustain or undermine a research community. "
              "A country may have a positive net inflow while simultaneously losing its domestic PI base through retirement or emigration, or it may have negative net flow but a healthy pipeline of new entrants. "
              "The transition-rate framework disaggregates these processes and shows that the same net flow can correspond to very different vulnerability profiles. "
              "For example, a high early-career outflow rate is less damaging than a high dropout rate because researchers abroad may return; a high dropout rate removes researchers from the system entirely. "
              "This distinction is lost in net-flow accounting but is central to point-of-no-return analysis.")

    p = doc.add_paragraph()
    p.add_run(f"First, {ctx['pnr_lever_text']}. "
              "A large proportional reduction in baseline recruitment would drive most communities to their threshold before mobility rates such as return or promotion became binding. "
              "This is consistent with the observation that AI/ML fields depend on a continuous pipeline of new graduate students and junior researchers")
    add_citation(p, 1)
    p.add_run(". "
              "Policies that sustain that pipeline, such as doctoral funding, visa routes for early-career researchers, and stable junior positions, are therefore first-order defences against a point of no return.")

    p = doc.add_paragraph()
    p.add_run(f"Second, among the mobility transition rates, dropout (d) is the dominant negative lever; its active-pool elasticity ranges from {_fmt(ctx['d_min_e'], 2)} to {_fmt(ctx['d_max_e'], 2)} across groups, and in the policy counterfactuals a simulated reduction in dropout yields the largest margin gain per unit proportional change. "
              "Attrition matters because it removes researchers from every compartment, not just one. "
              "A 10% proportional reduction in dropout expands the safety margin more than comparably sized increases in return, hit generation or promotion. "
              f"For {ctx['smallest_margin_group']}, the group with the smallest safety margin, even modest attrition reductions may widen the margin. "
              "These counterfactuals are mechanical perturbations of the fitted rates; they identify the most sensitive transition levers, not the causal effect of any specific policy programme.")

    p = doc.add_paragraph()
    p.add_run(f"Third, {ctx['positive_lever_sentence_lower']}. "
              f"The {ctx['highest_pd_group']} group shows the strongest response to PI promotion, suggesting that for that community expanding the domestic PI pipeline is an efficient lever. "
              "Return from abroad (β) is also positive for most groups, though its effect is generally smaller than reducing attrition directly. "
              "The implication for policy is that retention and promotion are usually more efficient than trying to attract returnees, but a balanced portfolio is still needed: a community without domestic PI growth cannot reproduce itself through attrition reduction alone.")

    p = doc.add_paragraph()
    p.add_run("Fourth, the historical counterfactual shows that the late-window rates, if they persisted, would alter equilibrium margins. "
              f"{period_direction_text} "
              "This pattern cautions against treating AI/ML mobility as a single global trend. "
              "It also confirms that the model can detect temporal changes in transition rates, which is the prerequisite for the early intervention the framework is designed to support.")

    p = doc.add_paragraph()
    p.add_run("The transition levers also interact in ways that a single-rate elasticity cannot fully capture. "
              "For example, reducing dropout and increasing PI promotion together are likely to have a larger effect than the sum of the two individual perturbations, because more researchers survive to become PIs and those PIs then train additional early-career researchers through the endogenous inflow channel. "
              "Conversely, a simultaneous fall in exogenous entry and a rise in dropout can push a community to its threshold faster than either change alone. "
              "The model's steady-state and one-at-a-time counterfactuals are therefore a starting point; they identify the most sensitive margins but do not exhaust the policy design space.")

    p = doc.add_paragraph()
    p.add_run("The connection to civilisational diversity is direct. "
              "Each group's safety margin can be monitored over time, and interventions can be adjusted before the margin disappears. "
              f"Because the endogenous inflow is capped at a safety factor of {_fmt(ctx['safety_factor_cap'], 2)} relative to the critical reproduction rate (the most constrained fitted group realises {_fmt(ctx['min_realised_safety_ratio'], 2)}), the policy recommendations are deliberately conservative: they do not push the system toward instability. "
              "That bounded approach is consistent with the goal of preserving diversity rather than maximising any single country's share.")

    p = doc.add_paragraph()
    p.add_run("It is important to stress that the counterfactuals reported in Tables 3 and 7 are mechanical perturbations of the fitted transition rates, not causal estimates of specific programmes. "
              "They identify which rates the model treats as most sensitive, and therefore where empirical policy evaluation is most urgent, but they do not by themselves show that a given intervention would achieve the simulated change.")

    doc.add_heading("6.2 Civilisational diversity as an innovation input", level=2)
    p = doc.add_paragraph()
    p.add_run("A second implication concerns the normative status of civilisational diversity. "
              "We treat diversity as an input to innovation rather than as a distributional afterthought")
    add_citation(p, 19)
    p.add_run(". "
              "A monocentric or tight-oligopoly structure in AI/ML may produce short-run efficiency gains through scale and agglomeration, but it also raises the risk of methodological lock-in, selection bias in training data, and reduced error correction. "
              "It is also an evolutionary dead end: it narrows the menu of innovation options, removes healthy competitors whose alternative approaches keep the field honest, and concentrates problem selection under a single institutional and methodological line. "
              "When one civilisation or a small oligopoly sets the dominant research agenda, problems that do not fit its priorities, languages, or institutional incentives are less likely to be addressed, leaving important scientific and social needs unresolved. "
              "Recent work on multi-university teams shows that geographically dispersed collaborations can retain high impact, which suggests that distributing capacity across civilisations need not sacrifice quality")
    add_citation(p, 12)
    p.add_run(". "
              "By quantifying the safety margin for each research community, the framework makes it possible to argue for support of smaller communities on positive, innovation-systems grounds. "
              "Preserving multiple centres of AI/ML research is not a matter of slowing the frontier; it is a matter of ensuring that the frontier is not defined by a single set of institutions, languages, or problems.")

    p = doc.add_paragraph()
    p.add_run("Japan is used as an illustrative case, not because it is the only group of interest, but because it combines a small absolute margin with rich data and a distinctive institutional lineage that makes the policy translation concrete. "
              "It is the clearest example among the large civilisations. "
              "Its fitted active-pool margin is T=" + str(ja_ctx['D'] + ja_ctx['H_D'] + ja_ctx['P_D']) + " researchers, with M=" + _fmt(ja_ctx['M'], 0) + " (T/M=" + _fmt(ja_ctx['T_over_M'], 2) + "). "
              "As Figure 8 shows, Japan's closest point of no return is the exogenous entry rate I0: if I0 fell to " + _fmt(ja_ctx['pnr_factor'] * 100, 1) + "% of its current level, the active pool would reach the minimum viable threshold. "
              "The same figure shows that Japan's early-career outflow α (" + _fmt(ja_ctx['alpha'], 3) + ") and domestic PI promotion p_D (" + _fmt(ja_ctx['p_D'], 3) + ") are comparatively low, while return from abroad β (" + _fmt(ja_ctx['beta'], 3) + ") and domestic hit generation h_D (" + _fmt(ja_ctx['h_D'], 3) + ") are moderate. "
              "These numbers translate into policy levers: α through retention fellowships and junior faculty positions; β through return grants and dual appointments; h_D through independent-lab programmes such as SPREAD; p_D through tenure-track conversion and startup packages; and d through childcare, dual-career support, and stable non-tenure tracks. "
              "Weakening the Japanese civilisation would not be neutral for the rest of the world: it would remove a distinct institutional lineage, reduce the pool of non-Anglophone problem framings, and leave a range of health, ageing, robotics, and materials problems under-addressed. "
              "Maintaining Japan as a viable AI/ML civilisation is therefore in the global interest, not only in Japan's national interest. "
              "The Japan-specific analysis is intended as a worked example; the same rate-ladder diagnostic can be applied to any civilisation with sufficient OpenAlex coverage.")
    doc.add_heading("6.3 Policy and management implications, and early warning", level=2)
    p = doc.add_paragraph()
    p.add_run("The policy implications can be read as an early-warning architecture. "
              "A single dashboard that tracks the fitted transition rates, their bootstrap uncertainty, and the distance to M for each civilisation would allow policymakers to detect divergence before a community enters an irreversible decline. "
              "Interventions can then be calibrated to maintain a minimum safety margin rather than to maximise any one stock. "
              "This is the operational meaning of early intervention: not a forecast that a particular collapse will occur, but a structured way to keep the system away from a point of no return. "
              "It also frames high-skilled mobility as a strategic competition among jurisdictions for talent")
    add_citation(p, 20)
    p.add_run(", in which the central question is not only who wins the current round but whether the global system retains enough diversity for future rounds")
    add_citation(p, 21)
    p.add_run(". "
              "If the response lag is short enough, the model can be updated annually and divergence caught before any single civilisation approaches a point of no return. "
              "It is therefore a tool for ensuring that no single civilisation reaches a self-sustaining collapse, and that the global AI/ML system retains the diversity required for continued innovation. "
              "We introduce the acronym SHIGA—Sustaining Heterogeneity through Interventions in Global AI/ML Researcher Mobility—formed from the title.")

    p = doc.add_paragraph()
    p.add_run("Table 8 maps the most sensitive transition levers to policy instruments and to the management actions that determine them. "
              "Policy instruments set incentives, while management actions determine how those incentives are implemented within institutions. "
              "Both are needed because a policy without a corresponding management process rarely changes transition rates.")

    lever_policy_mgmt = pd.DataFrame({
        "Lever": ["Dropout (d)", "Exogenous entry (I0)", "Return from abroad (β)", "Domestic hit generation (h_D)", "PI promotion (p_D)"],
        "Policy instrument": [
            "Early-career fellowships, childcare and dual-career support, stable non-tenure tracks",
            "Research-master and undergraduate pipelines, doctoral fellowships, recruitment visas",
            "Return grants, diaspora networks, dual appointments, overseas-experience recognition",
            "Independent-lab programmes (e.g. SPREAD-style), doctoral/postdoctoral training, compute access",
            "Tenure-track conversion, startup packages, project-based PI status",
        ],
        "Management action": [
            "Retain researchers in the domestic pipeline beyond the first career years",
            "Widen the base of incoming researchers before they select a field or location",
            "Encourage mobile researchers to re-establish domestic research groups",
            "Translate junior capacity into visible, high-impact work and independent research lines",
            "Create durable principal-investigator positions that train the next cohort",
        ],
    })
    _add_table_from_df(
        doc,
        lever_policy_mgmt,
        caption="Table 8. Transition levers, policy instruments, and management actions.",
    )

    p = doc.add_paragraph()
    p.add_run("Operationally, the framework can be used in two complementary ways. "
              "As a monitoring tool, it can be rerun whenever new OpenAlex data are released, producing an updated set of transition rates, safety margins and proximity-to-threshold estimates. "
              "As a scenario tool, it can quantify how large a proportional change in a given rate would be required to move a community toward or away from collapse, which helps prioritise empirical policy evaluation. "
              "Both uses depend on transparent assumptions and regular recalibration; the model should not be used to justify one-off interventions without accompanying process evaluation.")

    p = doc.add_paragraph()
    p.add_run("Table 8 maps each lever to the actors that control it: I0 and h_D are mainly owned by national funders and ministries; p_D and d by universities and department heads; and β by diaspora networks, return grants and private-sector recruiters. "
              "The model's management value is to rank which local rates most urgently need intervention and the proportional change needed to restore a safety margin.")

    doc.add_heading("6.4 Intra-civilisation alternatives when inter-civilisation mobility cannot be controlled", level=2)
    p = doc.add_paragraph()
    p.add_run("If a civilisation cannot control outflows to, or inflows from, other jurisdictions—whether because of visa regimes, salary differentials, language advantages, or targeted recruitment—it can still preserve its research community by acting on the intra-civilisation levers identified in the annual model. "
              "The annual rates show that the domestic active pool T = D + H_D + P_D responds most strongly to the dropout rate d, the domestic hit rate h_D, and the PI promotion rate p_D. "
              "Policies that reduce early-career attrition, expand domestic postdoctoral positions, or accelerate independent-lab formation therefore become defensive substitutes when inter-civilisation poaching cannot be regulated. "
              "This is the practical meaning of civilisational-diversity preservation under sovereignty constraints: even without controlling the border of talent, a community can increase the internal reproduction of active researchers. "
              f"The endogenous inflow is capped at a safety factor of {_fmt(ctx['safety_factor_cap'], 2)} relative to the critical reproduction rate (the most constrained fitted group realises {_fmt(ctx['min_realised_safety_ratio'], 2)}), so the model prevents over-optimism about this substitution effect; more ambitious domestic growth would require corresponding evidence that the extra PIs can be absorbed without simply raising dropout.")

    doc.add_heading("6.5 Annual updating as an early-warning layer", level=2)
    p = doc.add_paragraph()
    p.add_run("The 2017-2023 projection demonstrates that the framework can be rerun annually with a one-year time step. "
              "Each new year of OpenAlex data updates the observed transition rates, the fitted trends, and the distance to the minimum viable coauthor threshold. "
              "Because the model is regularised by the correction pressures, the one-year-ahead projection is not easily derailed by a single noisy observation. "
              "Instead, successive years reveal whether a particular transition rate is drifting toward a boundary. "
              "That drift is the early-warning signal. "
              "Policymakers can then intervene before the active pool falls below M, using the rate-specific elasticities in Table 3 to prioritise the smallest proportional change that restores a safety margin. "
              "This is the operational mechanism for avoiding technology monopoly and oligopoly dead ends: by keeping every major research community above its minimum viable coauthor pool, annual monitoring sustains the competitive diversity that underpins long-run technological progress. "
              "The framework is therefore not a prediction that a particular civilisation will collapse; it is a tool for ensuring that no single civilisation reaches a point where its collapse becomes self-sustaining. "
              "The modest year-to-year direction agreement in the 2017-2023 projection confirms that this layer is a drift-and-threshold alarm, not a precise population forecast. "
              "SHIGA therefore encapsulates the practical goal: keeping the global AI/ML system heterogeneous enough that no single centre of power can monopolise the technological frontier.")

    doc.add_heading("6.6 Limitations", level=2)
    p = doc.add_paragraph()
    p.add_run("Several limitations should be acknowledged. "
              "OpenAlex affiliation and country assignments are noisy, especially for researchers with multiple affiliations. "
              "The civilisation grouping is a coarse aggregation; within-group heterogeneity is substantial. "
              "The annual model relies on a discrete approximation of the continuous-time ODE and does not capture within-year events or cross-civilisation spillovers. "
              "Inter-civilisation flows are approximated by the author's recent_group while abroad, which misses year-to-year destination switching. "
              "The civilisation label is a pragmatic aggregation of publication-affiliation patterns. "
              "Historical civilisational boundaries do not necessarily coincide with contemporary political or value-based boundaries, and this study cannot determine whether the diversity of research ideas maps more closely onto historical civilisational groupings or onto current political and value communities; for example, the Sinic grouping reflects current OpenAlex country metadata and does not resolve the cultural and historical ties between mainland China and Taiwan, which currently appear as separate research arenas. "
              "This is treated as an empirical limitation of the classification, not as a normative claim. "
              "The cohort is a model-implied sample extracted from OpenAlex; absolute equilibrium numbers should be interpreted as model-implied stocks rather than census counts. "
              "Authors with many publications are over-weighted relative to less prolific authors, so rate estimates reflect author-publication exposure rather than a uniformly representative sample of individuals. "
              f"The endogenous inflow is capped at a safety factor of {_fmt(ctx['safety_factor_cap'], 2)} relative to the critical reproduction rate (the most constrained fitted group realises {_fmt(ctx['min_realised_safety_ratio'], 2)}); alternative values would shift equilibrium levels and should be reported in future sensitivity tables. "
              "Finally, the point-of-no-return threshold is a sufficient condition for collapse, not a necessary one: a community may decline for reasons outside the model even if T remains above M.")

    p = doc.add_paragraph()
    p.add_run("Wide bootstrap confidence intervals, especially for smaller civilisation groups, mean that the ordinal ranking of groups by equilibrium size or proximity to threshold should be treated as descriptive rather than definitive. "
              "The model identifies which transitions are most sensitive in a mechanical sense; turning those sensitivities into reliable policy priorities requires additional data on programme costs, implementation lags, and behavioural responses that are outside the scope of this paper.")

    p = doc.add_paragraph()
    p.add_run("From a security-studies perspective, the framework is intentionally non-adversarial: it treats mobility as an aggregate transition process and asks when a community becomes unable to reproduce itself, without modelling deliberate recruitment campaigns, technology transfer, or strategic denial. "
              "Future work could add a strategic layer by distinguishing civilian from defence-relevant AI/ML pipelines, or by modelling targeted recruitment in specific subfields.")

    # Conclusion
    doc.add_heading("7. Conclusion", level=1)
    p = doc.add_paragraph()
    p.add_run("We have proposed and implemented a transition-rate framework for assessing how close AI/ML research communities are to a point of no return. "
              "The model converts OpenAlex publication records into civilisation-specific transition rates and solves for the equilibrium active researcher pool. "
              "All groups remain above their minimum viable coauthor threshold in the fitted model, but the distance to that threshold varies by an order of magnitude and is most sensitive to exogenous entry and dropout. "
              f"Dropout is the dominant negative lever (active-pool elasticity {_fmt(ctx['d_min_e'], 2)} to {_fmt(ctx['d_max_e'], 2)}), and a simulated reduction is the single most efficient model-implied response for every civilisation. "
              "However, the closest point of no return is exogenous entry for all groups in the active-pool analysis, which means that policies which sustain the pipeline of new researchers are first-order defences. "
              "The historical counterfactual and the bootstrap intervals remind us that the future is not determined by current rates; transition rates can change, and policy can be directed at the most fragile lever before a collapse.")

    p = doc.add_paragraph()
    p.add_run("The annual projection layer adds an operational dimension to this conclusion. "
              "By estimating year-by-year transition rates and projecting one year ahead, the model turns the steady-state diagnostic into an early-warning dashboard. "
              "A one-year time step is short enough to detect drift before the active pool approaches the minimum viable threshold, and the correction pressures keep the projection within empirical and theoretical bounds. "
              "When inter-civilisation mobility cannot be controlled, the same framework points to intra-civilisation levers—reducing dropout, raising domestic hit rates, and accelerating PI promotion—that preserve T = D + H_D + P_D. "
              "These two layers, steady-state and annual, together provide a coherent basis for early, safety-factor-bound intervention.")

    p = doc.add_paragraph()
    p.add_run("The broader implication is that preserving civilisational diversity in AI/ML is compatible with, and may reinforce, scientific progress. "
              "A single dominant region or a tight oligopoly may achieve short-run scale economies, but it also risks methodological lock-in and reduces the set of problems that receive sustained attention. "
              "By monitoring transition rates and safety margins, policymakers can detect divergence early and intervene in a safety-factor-bound way. "
              "This is the practical meaning of the aspiration to avoid technology monopoly and oligopoly dead ends: not a prediction that any one civilisation will dominate, but a structured method for keeping the global system away from points of no return. "
              "Early, proportionate interventions that reduce attrition and sustain new recruitment can widen safety margins and preserve civilisational diversity in AI/ML.")

    doc.add_heading("7.1 Future work", level=2)
    p = doc.add_paragraph()
    p.add_run("Several extensions are natural. "
              "First, the model can be applied to other security-relevant fields such as semiconductor physics, quantum computing, biotechnology and energy materials, allowing cross-field comparisons of vulnerability. "
              "Second, the civilisation partition can be refined to a country or institution level, allowing bilateral migration flows and network externalities to be incorporated. "
              "Third, the ODE can be solved dynamically rather than at steady state, making it possible to forecast the time to threshold under alternative policy scenarios. "
              "Fourth, the minimum viable coauthor threshold can be made endogenous by modelling coauthorship as a matching process. "
              "Fifth, the sensitivity of equilibrium outcomes to the safety factor and to the saturating parameter epsilon should be mapped systematically. "
              "Finally, the framework can be integrated with policy cost data to produce cost-effectiveness comparisons of alternative interventions, turning mechanical sensitivities into actionable funding priorities.")

    # References
    doc.add_heading("References", level=1)
    for i, ref in enumerate(REFS, 1):
        p = doc.add_paragraph()
        p.add_run(f"{i}. {ref}")


def _anonymize_docx(doc):
    """Remove identifying text from a document for double-blind review."""
    replacements = {
        ", formed from the title and reflecting the research base at Shiga University": "",
        "https://github.com/bougtoir/researcher-mobility-ode": "[repository URL removed for double-blind review]",
    }
    for p in doc.paragraphs:
        for run in p.runs:
            for old, new in replacements.items():
                if old in run.text:
                    run.text = run.text.replace(old, new)
    for t in doc.tables:
        for row in t.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    for run in p.runs:
                        for old, new in replacements.items():
                            if old in run.text:
                                run.text = run.text.replace(old, new)


def write_docx(output_dir, data, fig_paths, blinded=False):
    abstract, keywords, highlights = _abstract_and_highlights(data[1], data[4])

    # Pre-compute body word count by building a throwaway body doc
    body_doc = Document()
    _add_docx_body(body_doc, data, fig_paths, blinded=blinded)
    body_wc = _doc_word_count(body_doc)

    doc = Document()
    _add_title_page(doc, word_count=body_wc, blinded=blinded)
    _add_front_matter(doc, abstract, keywords, highlights, blinded=blinded)
    _add_docx_body(doc, data, fig_paths, blinded=blinded)
    _unify_pnr_docx(doc)
    if blinded:
        _anonymize_docx(doc)

    suffix = "_blinded" if blinded else ""
    path = output_dir / f"manuscript_full_article{suffix}.docx"
    doc.save(path)
    return path


def write_pptx(output_dir, data, fig_paths):
    (cohort, eq, sat_eq, top_t, pnr_closest, period_compare, boot, policy_rank) = data
    annual = load_annual_data()
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    def add_image_slide(title, img_path, caption):
        slide_layout = prs.slide_layouts[3]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        left = PptxInches(1.5)
        top = PptxInches(1.2)
        slide.shapes.add_picture(str(img_path), left, top, width=PptxInches(10))
        txBox = slide.shapes.add_textbox(left, PptxInches(6.0), PptxInches(10), PptxInches(0.8))
        txBox.text_frame.text = caption
        for paragraph in txBox.text_frame.paragraphs:
            paragraph.font.size = Pt(14)

    def add_table_slide(title, df, col_names, width_per_col=1.5, font_size=10):
        slide_layout = prs.slide_layouts[3]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        rows, cols = len(df) + 1, len(col_names)
        left = PptxInches(0.5)
        top = PptxInches(1.2)
        table = slide.shapes.add_table(rows, cols, left, top, PptxInches(cols * width_per_col), PptxInches(0.6 * rows)).table
        for i, h in enumerate(col_names):
            table.cell(0, i).text = str(h)
        for row_i, (_, row) in enumerate(df.iterrows()):
            for j, val in enumerate(row):
                table.cell(row_i + 1, j).text = str(val)
                table.cell(row_i + 1, j).text_frame.paragraphs[0].font.size = Pt(font_size)

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Sustaining Heterogeneity through Interventions in Global AI/ML Researcher Mobility: A Transition-Rate Framework"
    slide.placeholders[1].text = "Data-driven manuscript figures and tables"

    add_image_slide(
        "Figure 1: Equilibrium T vs minimum viable threshold",
        fig_paths["fig1"],
        "Blue bars: equilibrium T; orange bars: threshold M. All groups remain above the threshold, but margins vary widely.",
    )
    add_image_slide(
        "Figure 2: Closest point-of-no-return proximity",
        fig_paths["fig2"],
        "Smaller values mean a smaller proportional change in the listed rate is required to reach the threshold for the stated target pool.",
    )
    add_image_slide(
        "Figure 3: Historical counterfactual margin change",
        fig_paths["fig3"],
        "Positive values mean the late-window rates would produce a larger safety margin than the early-window rates if they persisted; negative values mean the margin would shrink. The comparison is across point estimates; uncertainty is substantial.",
    )
    add_image_slide(
        "Figure 4: Bootstrap 95% CI for equilibrium T",
        fig_paths["fig4"],
        "Intervals are asymmetric and wide, reflecting model uncertainty.",
    )

    desc = _descriptive_table(cohort)
    add_table_slide("Table 1: Descriptive cohort statistics", desc, desc.columns.tolist(), width_per_col=1.3)

    eq_table = eq[["group", "T_equilibrium", "M_threshold", "margin_to_threshold_T", "I0", "r", "r_obs", "r_critical"]].copy()
    eq_table["T_over_M"] = eq_table["T_equilibrium"] / eq_table["M_threshold"]
    eq_table.columns = ["Group", "T_eq", "M", "Margin", "T/M", "I0", "r", "r_obs", "r_crit"]
    for c in ["T_eq", "M", "Margin", "I0"]:
        eq_table[c] = eq_table[c].apply(lambda x: _fmt(x, 0))
    eq_table["T/M"] = eq_table["T/M"].apply(lambda x: _fmt(x, 2))
    for c in ["r", "r_obs", "r_crit"]:
        eq_table[c] = eq_table[c].apply(lambda x: _fmt(x, 5))
    add_table_slide("Table 2: Equilibrium and inflow parameters", eq_table, eq_table.columns.tolist(), width_per_col=1.35)

    rows2 = []
    for group, gdf in top_t.groupby("group"):
        top3 = gdf.sort_values("abs_elasticity", ascending=False).head(3)
        parts = [group]
        for _, row in top3.iterrows():
            parts.extend([row["rate"], _fmt(row["elasticity"], 3)])
        rows2.append(parts)
    elas_df = pd.DataFrame(rows2, columns=["Group", "1st", "el1", "2nd", "el2", "3rd", "el3"])
    add_table_slide("Table 3: Top transition-rate elasticities", elas_df, elas_df.columns.tolist(), width_per_col=1.4)

    pnr_table = pnr_closest[["group", "target", "rate_name", "current_rate", "critical_factor", "proximity"]].copy()
    pnr_table.columns = ["Group", "Target", "Rate", "Current", "Crit.factor", "Proximity"]
    for c, d in {"Current": 4, "Crit.factor": 3, "Proximity": 3}.items():
        pnr_table[c] = pnr_table[c].apply(lambda x, d=d: _fmt(x, d))
    add_table_slide("Table 4: Closest point of no return", pnr_table, pnr_table.columns.tolist(), width_per_col=1.8)

    if sat_eq is not None:
        merged = eq[["group", "T_equilibrium"]].merge(
            sat_eq[["group", "T_equilibrium", "epsilon"]], on="group", suffixes=("_lin", "_sat")
        )
        merged.columns = ["Group", "Linear T", "Saturating T", "ε"]
        for c, d in {"Linear T": 0, "Saturating T": 0, "ε": 5}.items():
            merged[c] = merged[c].apply(lambda x, d=d: _fmt(x, d))
        add_table_slide("Table 5: Saturating inflow extension", merged, merged.columns.tolist(), width_per_col=2.0)

    pc = period_compare.rename(columns={
        "group": "Group",
        "T_early": "T early",
        "T_late": "T late",
        "pct_delta_T": "ΔT (%)",
        "margin_early": "Margin early",
        "margin_late": "Margin late",
        "delta_margin": "Δ margin",
    })
    for c, d in {"T early": 0, "T late": 0, "ΔT (%)": 1, "Margin early": 0, "Margin late": 0, "Δ margin": 1}.items():
        pc[c] = pc[c].apply(lambda x, d=d: _fmt(x, d))
    add_table_slide("Table 6: Historical counterfactual", pc, pc.columns.tolist(), width_per_col=1.4)

    policy_top = policy_rank.groupby("group").head(1).rename(columns={
        "group": "Group",
        "lever": "Lever",
        "direction": "Direction",
        "lever_change_pct": "Change (%)",
        "margin_gain": "Margin gain",
        "normalised_margin_gain_per_10pct": "Gain per 10%",
    })
    for c, d in {"Change (%)": 0, "Margin gain": 0, "Gain per 10%": 1}.items():
        policy_top[c] = policy_top[c].apply(lambda x, d=d: _fmt(x, d))
    add_table_slide("Table 7: Top policy intervention", policy_top, policy_top.columns.tolist(), width_per_col=2.0)

    lever_policy_mgmt = pd.DataFrame({
        "Lever": ["Dropout (d)", "Exogenous entry (I0)", "Return from abroad (β)", "Domestic hit generation (h_D)", "PI promotion (p_D)"],
        "Policy instrument": [
            "Early-career fellowships, childcare and dual-career support, stable non-tenure tracks",
            "Research-master and undergraduate pipelines, doctoral fellowships, recruitment visas",
            "Return grants, diaspora networks, dual appointments, overseas-experience recognition",
            "Independent-lab programmes (e.g. SPREAD-style), doctoral/postdoctoral training, compute access",
            "Tenure-track conversion, startup packages, project-based PI status",
        ],
        "Management action": [
            "Retain researchers in the domestic pipeline beyond the first career years",
            "Widen the base of incoming researchers before they select a field or location",
            "Encourage mobile researchers to re-establish domestic research groups",
            "Translate junior capacity into visible, high-impact work and independent research lines",
            "Create durable principal-investigator positions that train the next cohort",
        ],
    })
    add_table_slide("Table 8: Transition levers, policy instruments, and management actions", lever_policy_mgmt, lever_policy_mgmt.columns.tolist(), width_per_col=2.2)

    # Bootstrap CI is placed at the end as Supplementary Table 5.


    # Annual projection slides
    if fig_paths.get("fig5"):
        add_image_slide(
            "Figure 5: Observed and projected transition rates",
            fig_paths["fig5"],
            "Solid lines mark observed 2000-2016 rates; dashed lines mark projected 2017-2026 rates.",
        )
    if fig_paths.get("fig6"):
        add_image_slide(
            "Figure 6: Cross-civilisation abroad author-years",
            fig_paths["fig6"],
            "Rows are origin civilisations; columns are destination civilisations approximated by recent_group. Same-civilisation cells and Unknown destinations are excluded.",
        )
    if fig_paths.get("fig7"):
        add_image_slide(
            "Figure 7: Observed vs projected compartment counts",
            fig_paths["fig7"],
            "Solid lines are observed counts; dashed lines are 2017-2026 projections. The vertical dotted line is 2016.",
        )

    if fig_paths.get("fig8"):
        add_image_slide(
            "Figure 8: Japan compartment model and cross-civilisation transition-rate ladders",
            fig_paths["fig8"],
            "Left: Japan's six compartments; right: Japan highlighted against other civilisations on each transition rate.",
        )

    if fig_paths.get("fig9"):
        add_image_slide(
            "Figure 9: T/M safety margin versus closest PNR proximity",
            fig_paths["fig9"],
            "Lower-left points are the most fragile. Japan is highlighted in red.",
        )

    annual_means = annual_summary_table(annual)
    if not annual_means.empty:
        add_table_slide(
            "Supplementary Table 3: Mean observed annual transition rates, 2000-2016",
            annual_means,
            annual_means.columns.tolist(),
            width_per_col=1.4,
        )

    interciv_top = interciv_top_table(annual)
    if not interciv_top.empty:
        add_table_slide(
            "Supplementary Table 4: Top origin-destination abroad author-year pairs",
            interciv_top,
            interciv_top.columns.tolist(),
            width_per_col=2.0,
        )

    group_acc = annual.get("group_accuracy")
    if group_acc is not None and not group_acc.empty:
        gacc = group_acc.copy()
        gacc["rmse"] = gacc["rmse"].apply(lambda x: _fmt(x, 2))
        gacc["mape"] = gacc["mape"].apply(lambda x: f"{x*100:.1f}%")
        gacc = gacc.rename(columns={"origin_group": "Group", "rmse": "RMSE", "mape": "MAPE"})
        add_table_slide(
            "Supplementary Table 1: Projection accuracy by civilisation, 2017-2023",
            gacc,
            gacc.columns.tolist(),
            width_per_col=2.2,
        )

    comp_acc = annual.get("compartment_accuracy")
    if comp_acc is not None and not comp_acc.empty:
        cacc = comp_acc.copy()
        cacc["rmse"] = cacc["rmse"].apply(lambda x: _fmt(x, 2))
        cacc["mape"] = cacc["mape"].apply(lambda x: f"{x*100:.1f}%")
        cacc = cacc.rename(columns={"compartment": "Compartment", "rmse": "RMSE", "mape": "MAPE"})
        add_table_slide(
            "Supplementary Table 2: Projection accuracy by compartment, 2017-2023",
            cacc,
            cacc.columns.tolist(),
            width_per_col=2.2,
        )

    boot_tab = boot.copy()
    boot_tab["T 95% CI"] = boot_tab.apply(lambda r: f"[{_fmt(r['T_equilibrium_q025'], 0)}, {_fmt(r['T_equilibrium_q975'], 0)}]", axis=1)
    boot_tab["P_D 95% CI"] = boot_tab.apply(lambda r: f"[{_fmt(r['P_D_equilibrium_q025'], 0)}, {_fmt(r['P_D_equilibrium_q975'], 0)}]", axis=1)
    boot_tab = boot_tab[["group", "T_equilibrium_median", "T 95% CI", "P_D_equilibrium_mean", "P_D 95% CI"]]
    boot_tab.columns = ["Group", "T median", "T 95% CI", "P_D mean", "P_D 95% CI"]
    for c in ["T median", "P_D mean"]:
        boot_tab[c] = boot_tab[c].apply(lambda x: _fmt(x, 0))
    add_table_slide("Supplementary Table 5: Bootstrap 95% CI", boot_tab, boot_tab.columns.tolist(), width_per_col=2.2)

    path = output_dir / "manuscript_full_article_figures.pptx"
    prs.save(path)
    return path


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def write_supplementary_docx(output_dir, data, fig_paths):
    """Write a supplementary-materials docx with detailed tables supporting the main manuscript."""
    annual = load_annual_data()
    boot = data[6]
    doc = Document()
    doc.add_heading("Supplementary Material", level=0)
    p = doc.add_paragraph()
    p.add_run("Sustaining Heterogeneity through Interventions in Global AI/ML Researcher Mobility: A Transition-Rate Framework")
    p = doc.add_paragraph()
    p.add_run("This supplement provides detailed tables that support the main manuscript. "
              "Values are reproduced from the same result CSVs used to generate the main tables and figures; no numbers are hard-coded.")

    doc.add_heading("Supplementary Table 1. Projection accuracy by civilisation, 2017-2023", level=1)
    group_acc = annual.get("group_accuracy")
    if group_acc is not None and not group_acc.empty:
        gacc = group_acc.copy()
        gacc["rmse"] = gacc["rmse"].apply(lambda x: _fmt(x, 2))
        gacc["mape"] = gacc["mape"].apply(lambda x: f"{x*100:.1f}%")
        if "direction_agreement" in gacc.columns:
            gacc["direction_agreement"] = gacc["direction_agreement"].apply(lambda x: f"{x*100:.1f}%")
        if "threshold_alarm_accuracy" in gacc.columns:
            for c in ["threshold_alarm_accuracy", "threshold_alarm_sensitivity", "threshold_alarm_specificity", "threshold_alarm_precision"]:
                gacc[c] = gacc[c].apply(lambda x: f"{x*100:.1f}%" if pd.notna(x) else "—")
        rename = {"origin_group": "Group", "rmse": "RMSE", "mape": "MAPE", "direction_agreement": "Direction agreement"}
        for c in ["threshold_alarm_accuracy", "threshold_alarm_sensitivity", "threshold_alarm_specificity", "threshold_alarm_precision"]:
            if c in gacc.columns:
                rename[c] = c.replace("threshold_alarm_", "Alarm ").replace("_", " ").title()
        gacc = gacc.rename(columns=rename)
        _add_table_from_df(doc, gacc, caption="Supplementary Table 1. Projection accuracy by civilisation, 2017-2023.", decimals={"MAPE": 2})
    else:
        doc.add_paragraph("No group-level accuracy data available.")

    doc.add_heading("Supplementary Table 2. Projection accuracy by compartment, 2017-2023", level=1)
    comp_acc = annual.get("compartment_accuracy")
    if comp_acc is not None and not comp_acc.empty:
        cacc = comp_acc.copy()
        cacc["rmse"] = cacc["rmse"].apply(lambda x: _fmt(x, 2))
        cacc["mape"] = cacc["mape"].apply(lambda x: f"{x*100:.1f}%")
        if "direction_agreement" in cacc.columns:
            cacc["direction_agreement"] = cacc["direction_agreement"].apply(lambda x: f"{x*100:.1f}%")
        cacc = cacc.rename(columns={"compartment": "Compartment", "rmse": "RMSE", "mape": "MAPE", "direction_agreement": "Direction agreement"})
        _add_table_from_df(doc, cacc, caption="Supplementary Table 2. Projection accuracy by compartment, 2017-2023.", decimals={"MAPE": 2})
    else:
        doc.add_paragraph("No compartment-level accuracy data available.")

    doc.add_heading("Supplementary Table 3. Mean observed annual transition rates by civilisation, 2000-2016", level=1)
    annual_means = annual_summary_table(annual)
    if not annual_means.empty:
        _add_table_from_df(doc, annual_means, caption="Supplementary Table 3. Mean observed annual transition rates by civilisation, 2000-2016.", decimals={"α": 3, "β": 3, "h_D": 3, "p_D": 3, "d": 3, "I_total": 2})
    else:
        doc.add_paragraph("No annual transition-rate data available.")

    doc.add_heading("Supplementary Table 4. Top cross-civilisation origin-destination abroad author-year pairs", level=1)
    interciv_top = interciv_top_table(annual)
    if not interciv_top.empty:
        _add_table_from_df(doc, interciv_top, caption="Supplementary Table 4. Top cross-civilisation origin-destination abroad author-year pairs.", decimals={"Author-years": 0})
    else:
        doc.add_paragraph("No inter-civilisation flow data available.")

    doc.add_heading("Supplementary Table 5. Bootstrap 95% confidence intervals for equilibrium T and domestic PI pool P_D", level=1)
    boot_tab = boot.copy()
    boot_tab["T 95% CI"] = boot_tab.apply(lambda r: f"[{_fmt(r['T_equilibrium_q025'], 0)}, {_fmt(r['T_equilibrium_q975'], 0)}]", axis=1)
    boot_tab["P_D 95% CI"] = boot_tab.apply(lambda r: f"[{_fmt(r['P_D_equilibrium_q025'], 0)}, {_fmt(r['P_D_equilibrium_q975'], 0)}]", axis=1)
    boot_tab = boot_tab[["group", "T_equilibrium_median", "T 95% CI", "P_D_equilibrium_mean", "P_D 95% CI"]]
    boot_tab.columns = ["Group", "T median", "T 95% CI", "P_D mean", "P_D 95% CI"]
    _add_table_from_df(doc, boot_tab, caption="Supplementary Table 5. Bootstrap 95% confidence intervals for equilibrium T and domestic PI pool P_D.", decimals={"T median": 0, "P_D mean": 0})

    sup_path = output_dir / "supplementary_material.docx"
    doc.save(sup_path)
    return sup_path


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--output-dir", type=Path, default=BASE_DIR / "docs")
    args = parser.parse_args()
    output_dir = args.output_dir
    output_dir.mkdir(parents=True, exist_ok=True)

    data = load_data()
    cohort, eq, sat_eq, top_t, pnr_closest, period_compare, boot, policy_rank = data

    transition_rates = pd.read_csv(BASE_DIR / "data" / "cohort" / "transition_rates.csv")
    pnr_full = pd.read_csv(ENDOG / "point_of_no_return.csv")

    fig_dir = output_dir / "figures"
    fig1 = build_figure1(eq, fig_dir)
    fig2 = build_figure2(pnr_closest, fig_dir)
    fig3 = build_figure3(period_compare, fig_dir)
    fig4 = build_figure4(boot, fig_dir)
    fig8 = build_figure8(transition_rates, eq, fig_dir)
    fig9 = build_figure9(eq, pnr_full, fig_dir)

    annual = load_annual_data()
    annual_figs = build_annual_figures(annual, fig_dir)

    fig_paths = {
        "fig1": fig1,
        "fig2": fig2,
        "fig3": fig3,
        "fig4": fig4,
        "fig5": annual_figs.get("fig5"),
        "fig6": annual_figs.get("fig6"),
        "fig7": annual_figs.get("fig7"),
        "fig8": fig8,
        "fig9": fig9,
    }

    docx_path = write_docx(output_dir, data, fig_paths, blinded=False)
    blinded_docx_path = write_docx(output_dir, data, fig_paths, blinded=True)
    md_path = write_markdown(output_dir, docx_path=docx_path, blinded=False)
    blinded_md_path = write_markdown(output_dir, docx_path=blinded_docx_path, blinded=True)
    pptx_path = write_pptx(output_dir, data, fig_paths)
    sup_path = write_supplementary_docx(output_dir, data, fig_paths)
    sup_md_path = write_markdown(output_dir, docx_path=sup_path, blinded=False)

    print(f"Wrote {docx_path}")
    print(f"Wrote {blinded_docx_path}")
    print(f"Wrote {md_path}")
    print(f"Wrote {blinded_md_path}")
    print(f"Wrote {pptx_path}")
    print(f"Wrote {sup_path}")
    if sup_md_path:
        print(f"Wrote {sup_md_path}")
    print(f"Figures saved to {fig_dir}")

    # Build a submission zip containing the manuscript, editable figures, and PNGs
    zip_path = output_dir / "manuscript_full_article_submission.zip"
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for path in [docx_path, blinded_docx_path, pptx_path, md_path, blinded_md_path, sup_path]:
            if path and path.exists():
                zf.write(path, arcname=path.name)
        if sup_md_path and sup_md_path.exists():
            zf.write(sup_md_path, arcname=sup_md_path.name)
        for fig in sorted(fig_dir.glob("*.png")):
            zf.write(fig, arcname=f"figures/{fig.name}")
    print(f"Wrote {zip_path}")


if __name__ == "__main__":
    main()

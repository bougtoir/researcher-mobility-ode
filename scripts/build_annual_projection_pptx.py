#!/usr/bin/env python3
"""Generate an editable PPTX of the annual-rate/projection figures."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

REPO = Path(__file__).resolve().parents[1]
FIGURES = REPO / "docs" / "figures"
OUT = REPO / "docs" / "annual_rates_projection_figures.pptx"

slides = [
    {
        "title": "Figure 1. Annual observed and projected transition rates per civilisation",
        "caption": (
            "Observed rates (2000-2016) and projected rates (2017-2026) for "
            "the six ODE transition parameters, estimated per civilisation. "
            "Unreliable trends are replaced by the historical mean."
        ),
        "image": FIGURES / "annual_rates_by_group.png",
    },
    {
        "title": "Figure 2. Inter-civilisation abroad stock (origin x destination)",
        "caption": (
            "Total abroad author-years accumulated by origin (rows) and destination "
            "(columns) civilisation, 2000-2023. Destination is inferred from "
            "recent_group during years classified as abroad."
        ),
        "image": FIGURES / "annual_interciv_heatmap.png",
    },
    {
        "title": "Figure 3. Observed vs projected annual compartment counts",
        "caption": (
            "2017-2023 observed counts (solid) and projections from the 2000-2016 "
            "model (dashed). The vertical dotted line marks 2016, the end of the "
            "training window."
        ),
        "image": FIGURES / "annual_projection_vs_observed.png",
    },
]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank = prs.slide_layouts[6]
for s in slides:
    slide = prs.slides.add_slide(blank)
    # title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = s["title"]
    p.font.size = Pt(20)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    # image
    slide.shapes.add_picture(str(s["image"]), Inches(0.7), Inches(1.2), width=Inches(12.0))
    # caption
    cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.7))
    tf2 = cap_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = s["caption"]
    p2.font.size = Pt(14)
    p2.alignment = PP_ALIGN.LEFT

prs.save(OUT)
print(f"Wrote {OUT}")
